#!/usr/bin/env python3
"""
Unit III PowerPoint — Introduction to Infrastructure as Code (IaC)
Expanded deck: 4–5+ slides per syllabus topic with subtopics + real-life cases.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

NAVY = RGBColor(0x0A, 0x1B, 0x2E)
SKY = RGBColor(0x1A, 0x6F, 0xB5)
SKY_LIGHT = RGBColor(0x3D, 0x9B, 0xE0)
TEAL = RGBColor(0x0D, 0xA5, 0xA0)
CORAL = RGBColor(0xE8, 0x6A, 0x4C)
CLOUD = RGBColor(0xF4, 0xF8, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xE8, 0xF1, 0xF8)
INK = RGBColor(0x1C, 0x2B, 0x3A)
MUTED = RGBColor(0x5A, 0x6B, 0x7C)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
GREEN = RGBColor(0x1F, 0x8A, 0x5F)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
PAGE = 0


def add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_rounded(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def set_run(run, text, size=18, bold=False, color=INK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def tb(slide, left, top, width, height, lines, default_size=16, default_color=INK,
       align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, bold, color = item, default_size, False, default_color
        else:
            text = item[0]
            size = item[1] if len(item) > 1 else default_size
            bold = item[2] if len(item) > 2 else False
            color = item[3] if len(item) > 3 else default_color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(3)
        run = p.add_run()
        set_run(run, text, size=size, bold=bold, color=color)
    return box


def footer(slide):
    global PAGE
    PAGE += 1
    add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), NAVY)
    tb(slide, Inches(0.35), Inches(7.18), Inches(10.5), Inches(0.3),
       [("CC-II  ·  Unit III  ·  Infrastructure as Code (IaC)  ·  Parul University", 10, False, SOFT)])
    tb(slide, Inches(11.6), Inches(7.18), Inches(1.4), Inches(0.3),
       [(str(PAGE), 10, False, SOFT)], align=PP_ALIGN.RIGHT)


def header(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.08), TEAL)
    tb(slide, Inches(0.45), Inches(0.2), Inches(12.3), Inches(0.45),
       [(title, 24, True, WHITE)])
    if subtitle:
        tb(slide, Inches(0.45), Inches(0.68), Inches(12.3), Inches(0.32),
           [(subtitle, 12, False, SKY_LIGHT)])


def blank():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CLOUD)
    return slide


def card(slide, left, top, width, height, title, body_lines, accent=SKY):
    add_rounded(slide, left, top, width, height, WHITE)
    add_rect(slide, left, top, Inches(0.1), height, accent)
    tb(slide, left + Inches(0.22), top + Inches(0.12), width - Inches(0.35), Inches(0.35),
       [(title, 14, True, NAVY)])
    formatted = [(ln, 12, False, MUTED) if isinstance(ln, str) else ln for ln in body_lines]
    tb(slide, left + Inches(0.22), top + Inches(0.5), width - Inches(0.35), height - Inches(0.6),
       formatted)


def banner(slide, top, text, color=CORAL):
    add_rounded(slide, Inches(0.45), top, Inches(12.4), Inches(0.5), RGBColor(0xFF, 0xF4, 0xEC))
    tb(slide, Inches(0.65), top + Inches(0.08), Inches(12), Inches(0.35),
       [(f"Real life:  {text}", 12, True, color)])


def section_divider(section_no, title, bullets):
    s = blank()
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, TEAL)
    tb(s, Inches(0.8), Inches(2.0), Inches(11), Inches(0.4),
       [(f"TOPIC {section_no}", 14, True, TEAL)])
    tb(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.0),
       [(title, 34, True, WHITE)])
    for i, b in enumerate(bullets):
        tb(s, Inches(0.8), Inches(4.0 + i * 0.4), Inches(11.5), Inches(0.35),
           [(f"→  {b}", 14, False, SOFT)])
    footer(s)
    return s


# ═══════════════════════════════════════════════════════════════════════════
# OPENING
# ═══════════════════════════════════════════════════════════════════════════

s = blank()
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.15), TEAL)
tb(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.4),
   [("CLOUD COMPUTING – II  ·  UNIT III", 15, True, TEAL)])
tb(s, Inches(0.8), Inches(2.1), Inches(11.5), Inches(1.3),
   [("Infrastructure as Code", 42, True, WHITE), ("(IaC)", 42, True, WHITE)])
tb(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(1.2),
   [("Complete classroom deck  ·  Syllabus topics with subtopics", 14, False, SOFT),
    ("Overview  ·  Terraform / CloudFormation / Ansible", 14, False, SOFT),
    ("Configuration Management  ·  Workflow  ·  Versioning & Scaling", 14, False, SOFT)])
tb(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.4),
   [("FIT & CS  ·  Parul University  ·  Weightage 20%  ·  10 hours", 14, False, MUTED)])
PAGE += 1

s = blank()
header(s, "Learning Outcomes", "By the end of Unit III you should be able to…")
footer(s)
outcomes = [
    ("Define IaC", "Explain why code beats console click-ops for cloud infrastructure."),
    ("Compare tools", "Terraform vs CloudFormation vs Ansible — purpose, language, cloud, approach."),
    ("Split layers", "Provisioning creates resources; configuration management installs software."),
    ("Run the workflow", "Write → Git → Review → Plan → Apply → Verify — and why plan matters."),
    ("Scale safely", "Version infra in Git; parameterise capacity; promote Dev → Staging → Prod."),
    ("Link to DR", "Backups restore data; IaC rebuilds the system shape (Unit II connection)."),
]
for i, (t, d) in enumerate(outcomes):
    col = i % 2
    row = i // 2
    left = Inches(0.45 + col * 6.4)
    top = Inches(1.4 + row * 1.75)
    card(s, left, top, Inches(6.1), Inches(1.55), t, [d],
         accent=[SKY, TEAL, CORAL, GOLD, SKY, TEAL][i])

s = blank()
header(s, "Unit III Syllabus — Roadmap", "Five syllabus topics  ·  each expanded with subtopics")
footer(s)
topics = [
    ("1", "Overview of Infrastructure as Code (IaC)"),
    ("2", "Tools for IaC — Terraform, AWS CloudFormation, Ansible"),
    ("3", "Configuration Management & Automation"),
    ("4", "Writing & Managing Infrastructure as Code"),
    ("5", "Best Practices for Versioning & Scaling Infrastructure"),
]
for i, (n, t) in enumerate(topics):
    top = Inches(1.35 + i * 1.0)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.85), WHITE)
    add_rect(s, Inches(0.5), top, Inches(0.9), Inches(0.85), [SKY, TEAL, CORAL, GOLD, GREEN][i])
    tb(s, Inches(0.65), top + Inches(0.25), Inches(0.6), Inches(0.4),
       [(n, 22, True, WHITE)], align=PP_ALIGN.CENTER)
    tb(s, Inches(1.65), top + Inches(0.25), Inches(10.8), Inches(0.4),
       [(t, 16, True, NAVY)])

# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 1 — Overview of IaC
# ═══════════════════════════════════════════════════════════════════════════

section_divider(1, "Overview of Infrastructure as Code", [
    "Definition & what counts as infrastructure",
    "Problems with manual / console provisioning",
    "Benefits, imperative vs declarative, idempotence",
    "Drift, source of truth, link to Units I–II",
])

s = blank()
header(s, "1.1  What is Infrastructure as Code?", "Exam one-liner + full meaning")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(1.5), "Definition (write this)",
     ["Managing and provisioning computing infrastructure through machine-readable",
      "definition files — not by clicking a cloud console or typing one-off commands.",
      "Treat infrastructure like application code: versioned, reviewed, tested, automated."],
     accent=TEAL)
card(s, Inches(0.45), Inches(3.05), Inches(6.0), Inches(3.5), "Traditional way",
     ["Open GCP / AWS console",
      "Click Create VM / firewall / LB",
      "Hope you remember every setting",
      "Hard to repeat for Dev / Prod",
      "Slow after a disaster"],
     accent=CORAL)
card(s, Inches(6.7), Inches(3.05), Inches(6.15), Inches(3.5), "IaC way",
     ["Write a file describing desired resources",
      "Store it in Git",
      "Run Terraform / CloudFormation / …",
      "Same file recreates the stack anytime",
      "Plan → Apply → Verify"],
     accent=GREEN)

s = blank()
header(s, "1.2  What counts as “infrastructure”?", "More than just virtual machines")
footer(s)
items = [
    ("Networking", "VPC, subnets, routes, firewalls, NAT, DNS"),
    ("Compute", "VMs, MIG / ASG, GKE clusters, serverless hooks"),
    ("Storage & Data", "Disks, buckets, Cloud SQL / RDS, backup flags"),
    ("Security / IAM", "Service accounts, roles, bindings, keys (careful)"),
    ("Load balancing", "HTTP(S) LB, backends, health checks, CDN"),
    ("Observability", "Log sinks, alert policies (common in orgs)"),
]
for i, (t, d) in enumerate(items):
    col, row = i % 3, i // 3
    card(s, Inches(0.45 + col * 4.25), Inches(1.4 + row * 2.6), Inches(4.05), Inches(2.35),
         t, [d, "If a cloud API can create it,", "you can usually declare it in IaC."],
         accent=[SKY, TEAL, CORAL, GOLD, GREEN, SKY][i])

s = blank()
header(s, "1.3  Why Unit III matters", "Course story arc")
footer(s)
arcs = [
    ("Unit I", "What does good cloud design look like?", "Principles & Well-Architected"),
    ("Unit II", "What if it fails? How fast can we recover?", "DR · RTO · RPO · backups"),
    ("Unit III", "How do we create & recreate it reliably?", "IaC · tools · CM · workflow"),
]
for i, (u, q, hint) in enumerate(arcs):
    left = Inches(0.45 + i * 4.25)
    add_rounded(s, left, Inches(1.5), Inches(4.05), Inches(4.5), WHITE)
    add_rect(s, left, Inches(1.5), Inches(4.05), Inches(0.7), [SKY, CORAL, TEAL][i])
    tb(s, left + Inches(0.2), Inches(1.65), Inches(3.6), Inches(0.4), [(u, 18, True, WHITE)])
    tb(s, left + Inches(0.25), Inches(2.5), Inches(3.55), Inches(1.8),
       [(q, 15, True, NAVY)])
    tb(s, left + Inches(0.25), Inches(4.6), Inches(3.55), Inches(0.8),
       [(hint, 13, False, MUTED)])

s = blank()
header(s, "1.4  Snowflake servers & console problems", "Why manual provisioning fails at scale")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(1.35), "Snowflake server",
     ["A machine configured by hand for months — special packages, one-off fixes,",
      "undocumented tweaks. Nobody can recreate it. If it dies, the project dies with it."],
     accent=CORAL)
problems = [
    ("Click-ops only", "Steps live in memory / screenshots"),
    ("Environment drift", "Dev ≠ Prod firewalls & sizes"),
    ("Forgotten resources", "Orphan VMs → surprise bills"),
    ("Slow DR rebuild", "RTO balloons after zone loss"),
    ("No peer review", "Open ports go live unnoticed"),
    ("Scaling pain", "50 identical VMs by hand = chaos"),
]
for i, (t, d) in enumerate(problems):
    col, row = i % 3, i // 3
    card(s, Inches(0.45 + col * 4.25), Inches(2.95 + row * 1.9), Inches(4.05), Inches(1.7),
         t, [d], accent=[SKY, TEAL, GOLD, CORAL, GREEN, SKY][i])

s = blank()
header(s, "1.5  Benefits of IaC", "Expand each point in 7-mark answers")
footer(s)
bens = [
    ("Repeatability", "Same definition → same environment shape"),
    ("Version control", "Git shows who changed what and why"),
    ("Automation", "CI/CD plans and applies changes"),
    ("Fewer errors", "Less missed checkboxes & typos"),
    ("Living docs", "The repo IS the architecture"),
    ("Faster DR", "Re-apply stack from Git (Unit II)"),
    ("Collaboration", "PRs for infrastructure"),
    ("Auditability", "Reviewable desired state"),
]
for i, (t, d) in enumerate(bens):
    col, row = i % 4, i // 4
    card(s, Inches(0.4 + col * 3.2), Inches(1.4 + row * 2.6), Inches(3.05), Inches(2.35),
         t, [d], accent=[SKY, TEAL, CORAL, GOLD][i % 4])

s = blank()
header(s, "1.6  Imperative vs Declarative", "High-value exam distinction")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(4.0), "Imperative — HOW",
     ["Ordered steps / commands",
      "Example: gcloud compute instances create …",
      "Bash scripts that always “create”",
      "Re-run may duplicate unless you",
      "carefully code “if exists, skip”",
      "",
      "You are the orchestrator"],
     accent=CORAL)
card(s, Inches(6.8), Inches(1.35), Inches(6.05), Inches(4.0), "Declarative — WHAT",
     ["Declare desired end state",
      "Example: Terraform / CloudFormation",
      "Tool computes create/update/delete",
      "1st apply creates; 2nd apply is quiet",
      "Change count 2→5 → tool adds 3",
      "",
      "The tool converges to the target"],
     accent=TEAL)
banner(s, Inches(5.6), "Exam tip: Terraform & CloudFormation are declarative; Ansible playbooks are task-oriented.")

s = blank()
header(s, "1.7  Idempotence, drift & source of truth", "Must-know vocabulary")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(4.05), Inches(4.9), "Idempotent",
     ["Apply twice → same final state",
      "",
      "Good: package state=present",
      "Bad: echo >> file every run",
      "",
      "Safe re-runs after partial failure"],
     accent=TEAL)
card(s, Inches(4.7), Inches(1.35), Inches(4.05), Inches(4.9), "Drift",
     ["Actual cloud ≠ desired in Git",
      "",
      "Often from emergency console edits",
      "",
      "Fix: re-apply code, OR update",
      "code to include the real fix"],
     accent=CORAL)
card(s, Inches(8.95), Inches(1.35), Inches(3.9), Inches(4.9), "Source of truth",
     ["Git holds desired state",
      "",
      "Console is not the blueprint",
      "",
      "Emergency click → update IaC",
      "immediately afterward"],
     accent=GOLD)

s = blank()
header(s, "1.8  Provisioning vs configuration (preview)", "Full detail in Topic 3")
footer(s)
card(s, Inches(0.45), Inches(1.4), Inches(6.1), Inches(4.2), "Infrastructure provisioning",
     ["Create cloud resources via APIs",
      "VM, VPC, disk, LB, managed DB",
      "Tools: Terraform, CloudFormation",
      "",
      "“Build the house”"],
     accent=SKY)
card(s, Inches(6.8), Inches(1.4), Inches(6.05), Inches(4.2), "Configuration management",
     ["Install/configure software on servers",
      "Packages, app, users, services",
      "Tools: Ansible, startup scripts",
      "",
      "“Furnish the rooms”"],
     accent=TEAL)
banner(s, Inches(5.9), "Memory hook: Terraform builds the house; Ansible furnishes the rooms.")

# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 2 — Tools
# ═══════════════════════════════════════════════════════════════════════════

section_divider(2, "Tools for IaC", [
    "Master comparison table",
    "Terraform deep dive (HCL, state, plan/apply)",
    "AWS CloudFormation (templates, stacks, change sets)",
    "Ansible (playbooks, agentless SSH, CM focus)",
])

s = blank()
header(s, "2.1  Master comparison (memorise)", "Purpose · language · cloud · approach")
footer(s)
# table-like cards
headers_row = ["Tool", "Purpose", "Language", "Cloud", "Approach"]
rows = [
    ["Terraform", "Provision infra", "HCL", "Multi-cloud", "Declarative"],
    ["CloudFormation", "Provision AWS", "JSON / YAML", "AWS only", "Declarative"],
    ["Ansible", "Configure hosts", "YAML playbooks", "Multi / on-prem", "Task / procedural"],
]
y0 = Inches(1.4)
add_rounded(s, Inches(0.4), y0, Inches(12.5), Inches(0.55), NAVY)
widths = [2.2, 2.6, 2.5, 2.5, 2.7]
x = 0.5
for i, h in enumerate(headers_row):
    tb(s, Inches(x), y0 + Inches(0.12), Inches(widths[i]), Inches(0.35),
       [(h, 12, True, WHITE)])
    x += widths[i]
for r, row in enumerate(rows):
    top = Inches(2.15 + r * 1.15)
    add_rounded(s, Inches(0.4), top, Inches(12.5), Inches(1.0), WHITE)
    add_rect(s, Inches(0.4), top, Inches(0.12), Inches(1.0), [SKY, TEAL, CORAL][r])
    x = 0.55
    for i, cell in enumerate(row):
        tb(s, Inches(x), top + Inches(0.3), Inches(widths[i]), Inches(0.4),
           [(cell, 13, i == 0, NAVY if i == 0 else INK)])
        x += widths[i]
banner(s, Inches(5.9), "Extra viva: Deployment Manager (GCP), Pulumi, Azure Bicep — know names.")

s = blank()
header(s, "2.2  Terraform — what & why", "HashiCorp · multi-cloud · plan/apply culture")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(5.2), "Core building blocks",
     ["Provider — plugin for a cloud API",
      "Resource — thing you create/manage",
      "Variable — inputs (region, count)",
      "Output — values after apply (LB IP)",
      "Module — reusable package of resources",
      "State — maps code ↔ real cloud IDs"],
     accent=SKY)
card(s, Inches(6.8), Inches(1.35), Inches(6.05), Inches(5.2), "Essential commands",
     ["terraform init — providers / backend",
      "terraform fmt / validate",
      "terraform plan — preview (no change)",
      "terraform apply — execute plan",
      "terraform destroy — tear down",
      "terraform output — print outputs"],
     accent=TEAL)

s = blank()
header(s, "2.3  Terraform state & lifecycle", "Why teams use remote state + locking")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(1.4), "Lifecycle of a change",
     ["Edit .tf → fmt/validate → plan → human reviews plan → apply → verify health"],
     accent=TEAL)
card(s, Inches(0.45), Inches(3.0), Inches(6.1), Inches(3.5), "Local state (laptop)",
     ["OK for personal labs",
      "Risky for teams",
      "Lost laptop → lost mapping",
      "Two applies can race"],
     accent=CORAL)
card(s, Inches(6.8), Inches(3.0), Inches(6.05), Inches(3.5), "Remote state (GCS / S3)",
     ["Shared, durable source of mapping",
      "Locking = one apply at a time",
      "Restrict who can read state",
      "State may hold sensitive data"],
     accent=GREEN)

s = blank()
header(s, "2.4  Terraform — conceptual examples", "Exam: show you can sketch a resource")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(5.2), "Firewall (lab style)",
     ['resource "google_compute_firewall" "fw_fe" {',
      '  name    = "fw-fe"',
      '  network = "default"',
      '  allow { protocol = "tcp"',
      '          ports = ["8080"] }',
      '  target_tags = ["frontend"]',
      '}'],
     accent=SKY)
card(s, Inches(6.8), Inches(1.35), Inches(6.05), Inches(5.2), "Lab → resource map",
     ["VM → google_compute_instance",
      "Firewall → google_compute_firewall",
      "Template / MIG → instance_template /",
      "               instance_group_manager",
      "GKE → google_container_cluster",
      "Bucket → google_storage_bucket",
      "Cloud SQL → google_sql_database_instance"],
     accent=TEAL)

s = blank()
header(s, "2.5  AWS CloudFormation", "AWS-native IaC — templates, stacks, change sets")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(4.05), Inches(5.2), "Key terms",
     ["Template — YAML/JSON file",
      "Stack — deployed resource group",
      "Change set — preview updates",
      "Nested stack — modularity"],
     accent=SKY)
card(s, Inches(4.7), Inches(1.35), Inches(4.05), Inches(5.2), "Strengths",
     ["Deep AWS coverage",
      "First-party AWS support",
      "No separate state file you manage",
      "Fits AWS-only enterprises"],
     accent=TEAL)
card(s, Inches(8.95), Inches(1.35), Inches(3.9), Inches(5.2), "Limits",
     ["AWS only — not multi-cloud",
      "Huge templates get hard",
      "Multi-cloud orgs often",
      "standardise on Terraform"],
     accent=CORAL)

s = blank()
header(s, "2.6  CloudFormation — conceptual snippet", "S3 bucket in a stack")
footer(s)
card(s, Inches(0.45), Inches(1.4), Inches(12.4), Inches(4.5), "YAML idea",
     ["AWSTemplateFormatVersion: \"2010-09-09\"",
      "Resources:",
      "  WebBucket:",
      "    Type: AWS::S3::Bucket",
      "    Properties:",
      "      BucketName: my-cc2-demo-bucket",
      "      VersioningConfiguration:",
      "        Status: Enabled"],
     accent=GOLD)
banner(s, Inches(6.1), "Larger stacks also declare VPC, EC2, Security Groups, Load Balancers…")

s = blank()
header(s, "2.7  Ansible — configuration focus", "Agentless · YAML playbooks · often SSH push")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(5.2), "Vocabulary",
     ["Inventory — hosts / groups",
      "Playbook — YAML plays & tasks",
      "Module — apt, copy, service, user…",
      "Role — reusable task bundle",
      "Handler — e.g. restart nginx once"],
     accent=CORAL)
card(s, Inches(6.8), Inches(1.35), Inches(6.05), Inches(5.2), "Where it fits",
     ["Terraform / CFN create the VM",
      "Ansible installs packages & app",
      "Complements — not pure rivals",
      "",
      "Startup scripts = same idea",
      "expressed as Bash instead of YAML"],
     accent=TEAL)

s = blank()
header(s, "2.8  Ansible — conceptual playbook", "Idempotent server setup tasks")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(5.2), "Configure frontend hosts",
     ["- name: Install Node.js    →  apt: name=[nodejs,npm] state=present",
      "- name: Ensure app user    →  user: name=nodeapp state=present",
      "- name: Deploy app files   →  copy: src=./app/ dest=/opt/app/",
      "- name: Ensure service up  →  service: name=myapp state=started enabled=true",
      "",
      "Re-run safely: modules aim for idempotence (present / started)."],
     accent=SKY)

s = blank()
header(s, "2.9  Choosing a tool", "Decision guide for exams & labs")
footer(s)
choices = [
    ("Multi-cloud provisioning", "Terraform"),
    ("AWS-only native standard", "CloudFormation"),
    ("Configure OS / app on VMs", "Ansible"),
    ("GCP lab / Practical 7", "Terraform"),
    ("Create + configure together", "Terraform + Ansible"),
    ("K8s app deploy (Unit IV)", "YAML / Helm (declarative)"),
]
for i, (need, tool) in enumerate(choices):
    top = Inches(1.35 + i * 0.85)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.7), WHITE)
    tb(s, Inches(0.75), top + Inches(0.18), Inches(7.5), Inches(0.4), [(need, 14, False, INK)])
    tb(s, Inches(8.5), top + Inches(0.18), Inches(4.0), Inches(0.4),
       [(tool, 14, True, [SKY, TEAL, CORAL, GOLD, GREEN, SKY][i])])

# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 3 — Configuration Management
# ═══════════════════════════════════════════════════════════════════════════

section_divider(3, "Configuration Management & Automation", [
    "CM definition & goals",
    "Provisioning vs configuration (exam gold)",
    "Ansible / Chef / Puppet landscape",
    "8-step automated server setup + Fancy Store parallel",
])

s = blank()
header(s, "3.1  What is Configuration Management?", "Keep software & settings in desired state")
footer(s)
card(s, Inches(0.45), Inches(1.5), Inches(12.4), Inches(2.2), "Definition",
     ["Configuration Management maintains a system’s software and settings in a known,",
      "desired state — packages, files, services, users — usually through automation,",
      "not endless manual SSH."],
     accent=TEAL)
card(s, Inches(0.45), Inches(4.0), Inches(12.4), Inches(2.4), "Goal in one line",
     ["Every server of the same role looks the same — and stays that way.",
      "",
      "Without CM → drift, slow scale-out, fragile DR, security gaps on one host."],
     accent=CORAL)

s = blank()
header(s, "3.2  Provisioning vs configuration", "MiD1 Q.11 favourite — nail this table")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(5.2), "Infrastructure provisioning",
     ["Question: What cloud resources exist?",
      "",
      "Examples: VPC, VM, disk, firewall,",
      "LB, Cloud SQL instance",
      "",
      "Tools: Terraform, CloudFormation",
      "",
      "Timing: cloud control-plane APIs",
      "",
      "Blank VM created = provisioning done"],
     accent=SKY)
card(s, Inches(6.8), Inches(1.35), Inches(6.05), Inches(5.2), "Software configuration (CM)",
     ["Question: What runs on the machine?",
      "",
      "Examples: apt packages, Node.js,",
      "app files, nginx.conf, systemd, users",
      "",
      "Tools: Ansible, Chef, Puppet,",
      "startup scripts / cloud-init",
      "",
      "Timing: after OS/VM exists",
      "",
      "App serving traffic = CM done"],
     accent=TEAL)

s = blank()
header(s, "3.3  Why automation beats manual SSH", "Real operational pain")
footer(s)
pain = [
    ("Drift", "Friday hotfix on server-3 never copied to 1 & 2"),
    ("Slow scale-out", "Cannot confidently add 25 nodes before a sale"),
    ("Fragile DR", "Rebuild steps live in chat history"),
    ("Security gaps", "One host missing the hardening checklist"),
]
for i, (t, d) in enumerate(pain):
    top = Inches(1.4 + i * 1.25)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(1.05), WHITE)
    add_rect(s, Inches(0.5), top, Inches(0.12), Inches(1.05), [CORAL, GOLD, SKY, TEAL][i])
    tb(s, Inches(0.9), top + Inches(0.15), Inches(11.5), Inches(0.3), [(t, 15, True, NAVY)])
    tb(s, Inches(0.9), top + Inches(0.5), Inches(11.5), Inches(0.35), [(d, 13, False, MUTED)])

s = blank()
header(s, "3.4  CM tools & push vs pull", "Syllabus focus: Ansible")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(5.2), "Tool landscape",
     ["Ansible — SSH push; agentless; YAML",
      "Chef — recipes; often pull client",
      "Puppet — manifests; classic agents",
      "Salt — remote execution at scale",
      "cloud-init — first-boot from metadata",
      "",
      "For CC-II written exams: emphasise Ansible"],
     accent=CORAL)
card(s, Inches(6.8), Inches(1.35), Inches(6.05), Inches(5.2), "Push vs Pull",
     ["Push (common Ansible):",
      "Control node connects and pushes config",
      "",
      "Pull (some Chef/Puppet):",
      "Each node periodically fetches desired state",
      "",
      "Exam line: Ansible typically uses",
      "a push model over SSH and is agentless"],
     accent=TEAL)

s = blank()
header(s, "3.5  Eight-step automated server setup", "Memorise this ladder for Q.11")
footer(s)
steps = [
    "1  Update packages / install base OS tools",
    "2  Install runtime (Node.js, Java, Python…)",
    "3  Create dedicated non-root app user",
    "4  Pull code / artifacts (Git, GCS, Artifact Registry)",
    "5  Write config (ports, env; secrets from Secret Manager)",
    "6  Configure process manager (systemd / Supervisor)",
    "7  Start & enable service on boot",
    "8  Verify health (curl /health, check logs)",
]
for i, step in enumerate(steps):
    col, row = i % 2, i // 2
    top = Inches(1.35 + row * 1.25)
    left = Inches(0.45 + col * 6.4)
    add_rounded(s, left, top, Inches(6.15), Inches(1.05), WHITE)
    add_rect(s, left, top, Inches(0.12), Inches(1.05), [SKY, TEAL, CORAL, GOLD][row % 4])
    tb(s, left + Inches(0.3), top + Inches(0.3), Inches(5.6), Inches(0.45),
       [(step, 13, False, INK)])

s = blank()
header(s, "3.6  Fancy Store / lab parallel", "Startup script = CM in Bash form")
footer(s)
card(s, Inches(0.45), Inches(1.4), Inches(12.4), Inches(4.0), "Same pattern as Ansible",
     ["apt-get update & install tools",
      "→ install Node.js",
      "→ gsutil cp application from bucket",
      "→ configure Supervisor to keep process alive",
      "→ run as user nodeapp (not root)",
      "",
      "Ansible expresses the same intent as readable, reusable, idempotent tasks."],
     accent=TEAL)
banner(s, Inches(5.7), "Food-delivery fleets bake golden images so new VMs boot almost ready.")

s = blank()
header(s, "3.7  Images vs bootstrap", "Two strategies to get servers ready")
footer(s)
card(s, Inches(0.45), Inches(1.4), Inches(6.1), Inches(4.5), "A — Bootstrap at boot",
     ["Generic OS image",
      "Startup script / Ansible at boot",
      "Flexible; slower cold start",
      "More moving parts at runtime"],
     accent=SKY)
card(s, Inches(6.8), Inches(1.4), Inches(6.05), Inches(4.5), "B — Golden / baked image",
     ["Packer + Ansible bake packages in",
      "VM boots nearly ready",
      "Faster scale-out for peaks",
      "Requires image versioning"],
     accent=GREEN)

# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 4 — Writing & Managing IaC
# ═══════════════════════════════════════════════════════════════════════════

section_divider(4, "Writing & Managing Infrastructure as Code", [
    "Write → Git → Review → Plan → Apply → Verify",
    "Why plan-before-apply saves production",
    "Project layout, remote state, secrets",
    "PR checklist & CI/CD preview",
])

s = blank()
header(s, "4.1  The standard workflow", "Memorise exactly — MiD1 Q.12")
footer(s)
flow = [
    ("Write", "Author .tf / templates / playbooks"),
    ("Version\nControl", "Commit on a Git feature branch"),
    ("Review", "Pull request / peer review"),
    ("Plan", "Preview creates / changes / destroys"),
    ("Apply", "Execute approved changes"),
    ("Verify", "Health checks & smoke tests"),
]
for i, (t, d) in enumerate(flow):
    left = Inches(0.35 + i * 2.15)
    add_rounded(s, left, Inches(1.6), Inches(2.0), Inches(3.8), WHITE)
    add_rect(s, left, Inches(1.6), Inches(2.0), Inches(1.1), [SKY, TEAL, CORAL, GOLD, GREEN, SKY][i])
    tb(s, left + Inches(0.1), Inches(1.85), Inches(1.8), Inches(0.7),
       [(t, 14, True, WHITE)], align=PP_ALIGN.CENTER)
    tb(s, left + Inches(0.12), Inches(3.0), Inches(1.75), Inches(2.0),
       [(d, 12, False, MUTED)], align=PP_ALIGN.CENTER)
banner(s, Inches(5.8), "Skip any step and you invite drift, accidents, or unverified outages.")

s = blank()
header(s, "4.2  Why “Plan before Apply”?", "7-mark safety net")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(2.0), "Plan answers before impact",
     ["Will this destroy a database?",
      "Will this replace a VM (downtime)?",
      "Did I accidentally remove a firewall rule?"],
     accent=CORAL)
card(s, Inches(0.45), Inches(3.6), Inches(12.4), Inches(2.8), "Scare story (exam-friendly)",
     ["Engineer renames a Terraform resource without a state move.",
      "Plan shows: destroy old DB + create new empty DB.",
      "Without reading the plan, apply would be catastrophic.",
      "Plan-before-apply is production safety — not bureaucracy."],
     accent=GOLD)

s = blank()
header(s, "4.3  Typical Terraform project layout", "Modules = write once, reuse often")
footer(s)
layout = [
    "providers.tf / versions.tf — provider + required_version",
    "main.tf / network.tf — resources (split by concern)",
    "variables.tf — input declarations",
    "outputs.tf — LB IP, instance names…",
    "terraform.tfvars — env values (never commit secrets)",
    "modules/network · modules/compute — reusable blocks",
]
for i, line in enumerate(layout):
    top = Inches(1.35 + i * 0.85)
    add_rect(s, Inches(0.5), top + Inches(0.1), Inches(0.16), Inches(0.16),
             [SKY, TEAL, CORAL, GOLD, GREEN, SKY][i])
    tb(s, Inches(0.9), top, Inches(11.8), Inches(0.7), [(line, 14, False, INK)])

s = blank()
header(s, "4.4  Remote state, secrets & PR review", "Team safety essentials")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(4.05), Inches(5.2), "Remote state",
     ["GCS / S3 backend",
      "Locking enabled",
      "Tight IAM on bucket",
      "Separate state per env",
      "Survive laptop loss"],
     accent=SKY)
card(s, Inches(4.7), Inches(1.35), Inches(4.05), Inches(5.2), "Secrets — never in Git",
     ["No API keys / passwords",
      "No SA JSON key files",
      "Use Secret Manager",
      "Short-lived identities",
      "CI injects at runtime"],
     accent=CORAL)
card(s, Inches(8.95), Inches(1.35), Inches(3.9), Inches(5.2), "PR checklist",
     ["Least-privilege firewalls?",
      "Open SSH 0.0.0.0/0?",
      "DB deletion protection?",
      "Cost / machine size OK?",
      "Unexpected destroys?"],
     accent=TEAL)

s = blank()
header(s, "4.5  Stateful vs stateless changes", "Read the plan harder for databases")
footer(s)
card(s, Inches(0.45), Inches(1.5), Inches(6.1), Inches(4.3), "Stateless compute",
     ["Web frontends behind LB",
      "Rolling replace OK",
      "New template version → scale",
      "Easier to recreate"],
     accent=TEAL)
card(s, Inches(6.8), Inches(1.5), Inches(6.05), Inches(4.3), "Stateful data",
     ["Databases & unique disks",
      "Replace can mean data loss",
      "Use snapshots / PITR (Unit II)",
      "Extra caution on every plan"],
     accent=CORAL)

s = blank()
header(s, "4.6  IaC meets CI/CD (Unit V preview)", "Automation after human review")
footer(s)
steps2 = [
    "Developer opens PR with Terraform changes",
    "CI runs fmt · validate · plan (comment on PR)",
    "Human review + approval",
    "CD applies to Staging automatically",
    "Production apply may need a manual approval gate",
]
for i, line in enumerate(steps2):
    top = Inches(1.4 + i * 0.95)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.8), WHITE)
    add_rect(s, Inches(0.5), top, Inches(0.7), Inches(0.8), [SKY, TEAL, CORAL, GOLD, GREEN][i])
    tb(s, Inches(0.65), top + Inches(0.22), Inches(0.45), Inches(0.4),
       [(str(i + 1), 16, True, WHITE)], align=PP_ALIGN.CENTER)
    tb(s, Inches(1.45), top + Inches(0.22), Inches(11), Inches(0.4), [(line, 14, False, INK)])

# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 5 — Versioning & Scaling
# ═══════════════════════════════════════════════════════════════════════════

section_divider(5, "Versioning & Scaling Infrastructure", [
    "Git tags, pins, immutable artifacts",
    "Parameterise counts & autoscaling",
    "Dev → Staging → Prod promotion",
    "Checklists and anti-patterns",
])

s = blank()
header(s, "5.1  Versioning practices", "Immutable mindset")
footer(s)
practices = [
    ("Git as source of truth", "History, blame, collaboration"),
    ("Small clear PRs", "Reviewable intent"),
    ("Tag releases v1.0.0", "Correlate app ↔ infra"),
    ("Pin provider/modules", "No surprise breakage"),
    ("Immutable artifacts", "New image > mutate old"),
    ("Document breaking changes", "Team knows migrations"),
]
for i, (t, d) in enumerate(practices):
    col, row = i % 3, i // 3
    card(s, Inches(0.45 + col * 4.25), Inches(1.4 + row * 2.6), Inches(4.05), Inches(2.35),
         t, [d], accent=[SKY, TEAL, CORAL, GOLD, GREEN, SKY][i])

s = blank()
header(s, "5.2  Scaling with IaC — not copy-paste", "Capacity is a parameter")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(5.2), "Techniques",
     ["Parameterise instance_count / MIG size",
      "Encode autoscaling (CPU targets)",
      "Horizontal scale behind LB",
      "Vertical scale carefully via variables",
      "Multi-zone flags for HA (Unit II)",
      "Same module; different env sizes"],
     accent=SKY)
card(s, Inches(6.8), Inches(1.35), Inches(6.05), Inches(5.2), "Peak event example",
     ["Normal week: MIG size = 2",
      "Result day / sale: size = 12",
      "Same Terraform — new tfvars",
      "After peak: scale down (FinOps)",
      "",
      "Without IaC = panic click marathon"],
     accent=CORAL)

s = blank()
header(s, "5.3  Environment promotion", "Dev → Staging → Production")
footer(s)
envs = [
    ("Dev", "Break things cheaply", "Small counts · fast iteration"),
    ("Staging", "Production-like test", "Same modules · sale rehearsal"),
    ("Prod", "Extra approvals", "Scale params · careful plans"),
]
for i, (name, goal, hint) in enumerate(envs):
    left = Inches(0.5 + i * 4.2)
    add_rounded(s, left, Inches(1.6), Inches(3.95), Inches(4.2), WHITE)
    add_rect(s, left, Inches(1.6), Inches(3.95), Inches(0.9), [SKY, TEAL, CORAL][i])
    tb(s, left + Inches(0.2), Inches(1.8), Inches(3.5), Inches(0.5),
       [(name, 22, True, WHITE)], align=PP_ALIGN.CENTER)
    tb(s, left + Inches(0.25), Inches(2.9), Inches(3.4), Inches(1.2),
       [(goal, 16, True, NAVY)])
    tb(s, left + Inches(0.25), Inches(4.3), Inches(3.4), Inches(1.0),
       [(hint, 13, False, MUTED)])

s = blank()
header(s, "5.4  Operational checklist", "Printable habits")
footer(s)
checks = [
    "Git is source of truth — minimise long-lived console drift",
    "Always plan before apply in shared environments",
    "Remote state + locking for teams",
    "No secrets in Git",
    "Modules for reuse; variables for differences",
    "Labels/tags for ownership and cost",
    "Deletion protection + backups for data stores",
    "DR runbook includes “re-apply IaC” (Unit II)",
]
for i, c in enumerate(checks):
    top = Inches(1.3 + i * 0.65)
    tb(s, Inches(0.6), top, Inches(12), Inches(0.55),
       [(f"[  ]  {c}", 13, False, INK)])

s = blank()
header(s, "5.5  Anti-patterns to avoid", "Write these as “Avoid …” in exams")
footer(s)
antis = [
    "Click in console and never update code → permanent drift",
    "One gigantic root module nobody understands",
    "“God” credentials in plaintext committed tfvars",
    "terraform apply to prod from a laptop with no plan review",
    "IaC only on day 1, then forever SSH snowflake management",
    "Copy-pasting 50 near-identical blocks instead of count/modules",
]
for i, a in enumerate(antis):
    top = Inches(1.35 + i * 0.85)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.7), WHITE)
    add_rect(s, Inches(0.5), top, Inches(0.12), Inches(0.7), CORAL)
    tb(s, Inches(0.85), top + Inches(0.18), Inches(11.7), Inches(0.4), [(a, 13, False, INK)])

# ═══════════════════════════════════════════════════════════════════════════
# CASE STUDIES
# ═══════════════════════════════════════════════════════════════════════════

s = blank()
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, GOLD)
tb(s, Inches(0.8), Inches(2.3), Inches(11), Inches(0.4),
   [("CASE STUDIES", 14, True, GOLD)])
tb(s, Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.0),
   [("Real-life examples for 7-mark answers", 32, True, WHITE)])
tb(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.8),
   [("Quote one case + tool choice + workflow step for full marks", 14, False, SOFT)])
footer(s)

cases = [
    ("Campus result portal", "Same Terraform module; bump MIG size on result day; scale down after.",
     "Scale is a variable — not a click marathon."),
    ("Food delivery dinner peak", "Terraform autoscaling + golden image / Ansible so new VMs match.",
     "Provisioning + configuration both required."),
    ("UPI / payments (low RTO)", "DR runbook: restore data + re-apply IaC in another zone/region.",
     "IaC accelerates Disaster Recovery."),
    ("E-commerce festival sale", "Staging clones prod shape; PR catches open security group pre-sale.",
     "Review + plan prevent festival disasters."),
    ("Startup console drift", "Rewrite architecture into Terraform; ban long-lived click-ops.",
     "Tribal knowledge becomes onboarding."),
    ("DB destroy caught by plan", "Rename shows destroy/create; reviewer blocks; state move instead.",
     "Plan-before-apply saves production."),
]
for batch in (cases[:3], cases[3:]):
    s = blank()
    header(s, "Real-life cases", "Use these stories in long answers & viva")
    footer(s)
    for i, (title, approach, point) in enumerate(batch):
        top = Inches(1.35 + i * 1.8)
        add_rounded(s, Inches(0.45), top, Inches(12.4), Inches(1.6), WHITE)
        add_rect(s, Inches(0.45), top, Inches(0.12), Inches(1.6), [SKY, TEAL, CORAL][i])
        tb(s, Inches(0.8), top + Inches(0.15), Inches(11.8), Inches(0.35),
           [(title, 15, True, NAVY)])
        tb(s, Inches(0.8), top + Inches(0.55), Inches(11.8), Inches(0.4),
           [(approach, 12, False, INK)])
        tb(s, Inches(0.8), top + Inches(1.05), Inches(11.8), Inches(0.35),
           [(f"Teaching point: {point}", 12, True, MUTED)])

# ═══════════════════════════════════════════════════════════════════════════
# EXAM + CLOSE
# ═══════════════════════════════════════════════════════════════════════════

s = blank()
header(s, "Exam Rapid Revision — Sheet A", "Definitions & distinctions")
footer(s)
lines = [
    "IaC = infra as versioned machine-readable code (not console click-ops)",
    "Declarative = WHAT desired state · Imperative = HOW steps",
    "Idempotent = safe re-apply · Drift = actual ≠ desired in Git",
    "Provisioning creates cloud resources · CM configures software on hosts",
    "Terraform: multi-cloud, HCL, state, plan/apply",
    "CloudFormation: AWS-only, YAML/JSON, stacks & change sets",
    "Ansible: CM playbooks, YAML, agentless SSH push",
]
for i, line in enumerate(lines):
    top = Inches(1.3 + i * 0.75)
    add_rect(s, Inches(0.5), top + Inches(0.1), Inches(0.16), Inches(0.16),
             SKY if i % 2 == 0 else TEAL)
    tb(s, Inches(0.9), top, Inches(11.8), Inches(0.65), [(line, 13, False, INK)])

s = blank()
header(s, "Exam Rapid Revision — Sheet B", "Workflow · practices · must-say phrases")
footer(s)
lines = [
    "Workflow: Write → Git → Review → Plan → Apply → Verify",
    "Plan before apply prevents accidental DB destroy / costly mistakes",
    "CM 8 steps: packages → runtime → user → artifacts → config → PM → start → health",
    "Version: Git tags, pin versions, immutable images/templates",
    "Scale: parameterise counts, autoscaling, Dev→Staging→Prod",
    "Must-say: Git is source of truth · Plan before apply · TF provisions / Ansible configures",
    "Link Unit II: backups restore data; IaC rebuilds shape → lower RTO",
]
for i, line in enumerate(lines):
    top = Inches(1.3 + i * 0.75)
    add_rect(s, Inches(0.5), top + Inches(0.1), Inches(0.16), Inches(0.16),
             CORAL if i % 2 == 0 else GOLD)
    tb(s, Inches(0.9), top, Inches(11.8), Inches(0.65), [(line, 13, False, INK)])

s = blank()
header(s, "MiD1 Unit III map", "Q.9–Q.12 → which slides to revise")
footer(s)
qs = [
    ("Q.9  [7]", "What is IaC? Why over console? Benefits + one scenario", "Topics 1 + Cases"),
    ("Q.10 [7]", "Compare Terraform / CloudFormation / Ansible + example", "Topic 2"),
    ("Q.11 [7]", "CM vs provisioning · Ansible · server setup steps", "Topic 3"),
    ("Q.12 [7]", "Workflow · versioning/scaling · plan-before-apply", "Topics 4–5"),
]
for i, (q, prompt, where) in enumerate(qs):
    top = Inches(1.35 + i * 1.3)
    add_rounded(s, Inches(0.45), top, Inches(12.4), Inches(1.1), WHITE)
    tb(s, Inches(0.7), top + Inches(0.15), Inches(2.2), Inches(0.35), [(q, 14, True, NAVY)])
    tb(s, Inches(3.0), top + Inches(0.15), Inches(7.0), Inches(0.7), [(prompt, 13, False, INK)])
    tb(s, Inches(10.2), top + Inches(0.3), Inches(2.4), Inches(0.4),
       [(where, 12, True, TEAL)])

s = blank()
header(s, "Key Takeaways", "Ship infrastructure like software")
footer(s)
takes = [
    ("Code is the blueprint", "If it is not in Git, it is not the source of truth."),
    ("Plan before apply", "Preview destroys and replacements before production feels them."),
    ("Split the layers", "Terraform/CFN provision; Ansible configures — use both when needed."),
    ("Parameterise scale", "Result day / sale capacity is a variable — then scale back down."),
    ("IaC helps DR", "Restore data + re-apply infra code to cut RTO (Unit II link)."),
]
for i, (t, d) in enumerate(takes):
    top = Inches(1.3 + i * 1.05)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.9), WHITE)
    add_rect(s, Inches(0.5), top, Inches(0.12), Inches(0.9), [SKY, TEAL, CORAL, GOLD, GREEN][i])
    tb(s, Inches(0.9), top + Inches(0.12), Inches(11.5), Inches(0.3), [(t, 14, True, NAVY)])
    tb(s, Inches(0.9), top + Inches(0.45), Inches(11.5), Inches(0.35), [(d, 12, False, MUTED)])

s = blank()
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.15), TEAL)
tb(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.8), [("Thank you", 44, True, WHITE)])
tb(s, Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.5),
   [("Questions · Pick a case · Map it to a tool + workflow step", 16, False, SKY_LIGHT)])
tb(s, Inches(0.8), Inches(4.2), Inches(11.5), Inches(1.6),
   [("Unit III — Introduction to Infrastructure as Code (IaC)", 14, False, SOFT),
    ("Pair with: notes/unit-3/CC-II-Unit-III-Infrastructure-as-Code.txt", 14, False, SOFT),
    ("Lab link: Practical 7 — Terraform", 14, False, SOFT),
    ("FIT & CS, Parul University", 14, False, MUTED)])
PAGE += 1

out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "CC-II-Unit-III-Infrastructure-as-Code.pptx")
prs.save(out)
print(f"Saved {out}")
print(f"Slides: {len(prs.slides)}")
