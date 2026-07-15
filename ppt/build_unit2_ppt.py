#!/usr/bin/env python3
"""
Unit II PowerPoint — Disaster Recovery in Cloud
Expanded deck: 4–5+ slides per syllabus topic with subtopics.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

NAVY = RGBColor(0x0A, 0x1B, 0x2E)
SKY = RGBColor(0x1A, 0x6F, 0xB5)
SKY_LIGHT = RGBColor(0x3D, 0x9B, 0xE0)
TEAL = RGBColor(0x0D, 0xA5, 0xA0)
CORAL = RGBColor(0xE8, 0x6A, 0x4C)
CLOUD = RGBColor(0xF4, 0xF8, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xE8, 0xF1, 0xF8)
INK = RGBColor(0x1C, 0x2B, 0x3A)
MUTED = RGBColor(0x5A, 0x6B, 0x7C)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
GREEN = RGBColor(0x1F, 0x8A, 0x5F)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
PAGE = 0


def add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_rounded(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def set_run(run, text, size=18, bold=False, color=INK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def tb(slide, left, top, width, height, lines, default_size=16, default_color=INK,
       align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, bold, color = item, default_size, False, default_color
        else:
            text = item[0]
            size = item[1] if len(item) > 1 else default_size
            bold = item[2] if len(item) > 2 else False
            color = item[3] if len(item) > 3 else default_color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(3)
        run = p.add_run()
        set_run(run, text, size=size, bold=bold, color=color)
    return box


def footer(slide):
    global PAGE
    PAGE += 1
    add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), NAVY)
    tb(slide, Inches(0.35), Inches(7.18), Inches(10.5), Inches(0.3),
       [("CC-II  ·  Unit II  ·  Disaster Recovery in Cloud  ·  Parul University", 10, False, SOFT)])
    tb(slide, Inches(11.6), Inches(7.18), Inches(1.4), Inches(0.3),
       [(str(PAGE), 10, False, SOFT)], align=PP_ALIGN.RIGHT)


def header(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.08), TEAL)
    tb(slide, Inches(0.45), Inches(0.2), Inches(12.3), Inches(0.45),
       [(title, 24, True, WHITE)])
    if subtitle:
        tb(slide, Inches(0.45), Inches(0.68), Inches(12.3), Inches(0.32),
           [(subtitle, 12, False, SKY_LIGHT)])


def blank():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CLOUD)
    return slide


def card(slide, left, top, width, height, title, body_lines, accent=SKY):
    add_rounded(slide, left, top, width, height, WHITE)
    add_rect(slide, left, top, Inches(0.1), height, accent)
    tb(slide, left + Inches(0.22), top + Inches(0.12), width - Inches(0.35), Inches(0.35),
       [(title, 14, True, NAVY)])
    formatted = [(ln, 12, False, MUTED) if isinstance(ln, str) else ln for ln in body_lines]
    tb(slide, left + Inches(0.22), top + Inches(0.5), width - Inches(0.35), height - Inches(0.6),
       formatted)


def banner(slide, top, text, color=CORAL):
    add_rounded(slide, Inches(0.45), top, Inches(12.4), Inches(0.5), RGBColor(0xFF, 0xF4, 0xEC))
    tb(slide, Inches(0.65), top + Inches(0.08), Inches(12), Inches(0.35),
       [(f"Real life:  {text}", 12, True, color)])


def section_divider(section_no, title, bullets):
    """Full-bleed section opener."""
    s = blank()
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, TEAL)
    tb(s, Inches(0.8), Inches(2.0), Inches(11), Inches(0.4),
       [(f"TOPIC {section_no}", 14, True, TEAL)])
    tb(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.0),
       [(title, 36, True, WHITE)])
    for i, b in enumerate(bullets):
        tb(s, Inches(0.8), Inches(4.0 + i * 0.4), Inches(11.5), Inches(0.35),
           [(f"→  {b}", 14, False, SOFT)])
    footer(s)
    return s


def bullets_slide(title, subtitle, items, accent=SKY):
    s = blank()
    header(s, title, subtitle)
    footer(s)
    for i, item in enumerate(items):
        top = Inches(1.4 + i * 0.7)
        if top > Inches(6.5):
            break
        add_rect(s, Inches(0.5), top + Inches(0.08), Inches(0.16), Inches(0.16),
                 accent if i % 2 == 0 else TEAL)
        if isinstance(item, str):
            tb(s, Inches(0.85), top, Inches(11.9), Inches(0.55), [(item, 14, False, INK)])
        else:
            tb(s, Inches(0.85), top, Inches(11.9), Inches(0.55),
               [(item[0], 14, True, NAVY), (item[1], 12, False, MUTED)] if len(item) > 1 and False else [item])
    return s


def simple_bullets(title, subtitle, lines, note=None):
    s = blank()
    header(s, title, subtitle)
    footer(s)
    for i, line in enumerate(lines):
        top = Inches(1.35 + i * 0.65)
        add_rect(s, Inches(0.5), top + Inches(0.1), Inches(0.16), Inches(0.16),
                 [SKY, TEAL, CORAL, GOLD][i % 4])
        tb(s, Inches(0.9), top, Inches(11.8), Inches(0.55), [(line, 14, False, INK)])
    if note:
        banner(s, Inches(6.35) if len(lines) < 7 else Inches(6.4), note)
    return s


# ═══════════════════════════════════════════════════════════════════════════
# OPENING
# ═══════════════════════════════════════════════════════════════════════════

s = blank()
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.15), TEAL)
tb(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.4),
   [("CLOUD COMPUTING – II  ·  UNIT II", 15, True, TEAL)])
tb(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2),
   [("Disaster Recovery", 44, True, WHITE), ("in the Cloud", 44, True, WHITE)])
tb(s, Inches(0.8), Inches(4.4), Inches(11.5), Inches(1.0),
   [("Complete classroom deck  ·  Syllabus topics with subtopics", 14, False, SOFT),
    ("Backup & DR  ·  HA & Fault Tolerance  ·  RPO/RTO Strategies", 14, False, SOFT),
    ("Cloud Backups & Snapshots  ·  Case Studies", 14, False, SOFT)])
tb(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.4),
   [("FIT & CS  ·  Parul University", 14, False, MUTED)])
PAGE += 1  # title has no footer strip looking odd on navy — still count


s = blank()
header(s, "Learning Outcomes", "By the end of Unit II you should be able to…")
footer(s)
outcomes = [
    "Define Backup, Disaster, DR, Business Continuity, Failover, and Failback accurately",
    "Explain Full / Incremental / Snapshot backups and Backup vs Replication",
    "Differentiate High Availability, Fault Tolerance, Redundancy, and SPOF",
    "Calculate thinking for RTO & RPO and choose a DR strategy tier that fits them",
    "Map GCP tools (snapshots, GCS versioning, Cloud SQL PITR) to recovery playbooks",
    "Analyse case studies (lab + industry) and recommend an appropriate DR/HA design",
]
for i, o in enumerate(outcomes):
    top = Inches(1.4 + i * 0.85)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.7), WHITE)
    add_rect(s, Inches(0.5), top, Inches(0.12), Inches(0.7), [SKY, TEAL, CORAL, GOLD, SKY, TEAL][i])
    tb(s, Inches(0.9), top + Inches(0.18), Inches(11.6), Inches(0.4), [(o, 14, False, INK)])


s = blank()
header(s, "Unit II Syllabus — Roadmap", "Six syllabus topics  ·  each expanded with subtopics")
footer(s)
topics = [
    ("01", "Disaster Recovery in Cloud", "Intro · types of disasters · cloud DR advantage"),
    ("02", "Backup & DR Concepts", "Definitions · backup types · replication · 3-2-1"),
    ("03", "HA & Fault Tolerance", "HA · FT · SPOF · patterns · uptime"),
    ("04", "DR Strategies · RPO & RTO", "Metrics · ladder · each strategy deep dive"),
    ("05", "Cloud Backups & Snapshots", "GCE · GCS · Cloud SQL · GKE · practices"),
    ("06", "Case Studies", "Lab + industry DR implementations"),
]
for i, (num, title, sub) in enumerate(topics):
    col, row = i % 3, i // 3
    left = Inches(0.45 + col * 4.25)
    top = Inches(1.4 + row * 2.5)
    add_rounded(s, left, top, Inches(4.05), Inches(2.2), WHITE)
    add_rect(s, left, top, Inches(4.05), Inches(0.55), [SKY, TEAL, CORAL, SKY, TEAL, CORAL][i])
    tb(s, left + Inches(0.2), top + Inches(0.12), Inches(3.6), Inches(0.4),
       [(num + "  " + title, 13, True, WHITE)])
    tb(s, left + Inches(0.25), top + Inches(0.9), Inches(3.55), Inches(1.0),
       [(sub, 13, False, MUTED)])


# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 1 — Disaster Recovery in Cloud  (5+ slides)
# ═══════════════════════════════════════════════════════════════════════════

section_divider("01", "Disaster Recovery in Cloud", [
    "What DR means in a cloud environment",
    "What counts as a disaster",
    "Why Unit II matters after Unit I networking",
    "Cloud advantages and limits",
    "Shared responsibility for recovery",
])

s = blank()
header(s, "1.1  What is Disaster Recovery in Cloud?", "Policies + people + technology to restore operations after a disaster")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "Definition", [
    "Disaster Recovery (DR) is the organised capability",
    "to restore IT services after an event that makes",
    "systems unavailable or data unusable.",
    "",
    "It includes:",
    "• Documented policies & runbooks",
    "• Assigned people / roles",
    "• Technical tools (backups, replicas, failover)",
    "• Regular practice (game days / restore tests)",
], SKY)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "On-prem vs Cloud DR", [
    "Traditional on-prem:",
    "• Second physical datacenter",
    "• Tape trucks / slow rebuilds",
    "• Large idle hardware cost",
    "",
    "Cloud DR:",
    "• Spare capacity on demand",
    "• Snapshots & managed DB PITR",
    "• Global load balancing",
    "• Rebuild via Infrastructure as Code",
], TEAL)

s = blank()
header(s, "1.2  What Counts as a “Disaster”?", "Not only earthquakes — many day-to-day failures are DR events")
footer(s)
disasters = [
    ("Zone failure", "Power, cooling, or network in one datacenter location"),
    ("Region outage", "Broader geographic impact across a cloud region"),
    ("Cyberattack / ransomware", "Data encrypted or destroyed; access denied"),
    ("Accidental delete / bad deploy", "Human or pipeline change corrupts production"),
    ("Natural disaster", "Flood, fire, earthquake near a facility"),
    ("Dependency cascade", "DNS, identity, payment gateway, or SaaS dependency fails"),
]
for i, (t, d) in enumerate(disasters):
    col, row = i % 3, i // 3
    left = Inches(0.45 + col * 4.25)
    top = Inches(1.35 + row * 2.55)
    add_rounded(s, left, top, Inches(4.05), Inches(2.3), WHITE)
    add_rect(s, left, top, Inches(4.05), Inches(0.55), [CORAL, SKY, TEAL, GOLD, CORAL, SKY][i])
    tb(s, left + Inches(0.2), top + Inches(0.12), Inches(3.65), Inches(0.4),
       [(t, 14, True, WHITE)], align=PP_ALIGN.CENTER)
    tb(s, left + Inches(0.25), top + Inches(0.9), Inches(3.55), Inches(1.1),
       [(d, 13, False, MUTED)], align=PP_ALIGN.CENTER)

s = blank()
header(s, "1.3  Why Unit II Matters", "Unit I connects resources — Unit II asks what happens when they fail")
footer(s)
simple_bullets.__wrapped__ if False else None
qs = [
    "What happens when those networked resources fail?",
    "How much downtime is acceptable to the business? (RTO)",
    "How much data can we afford to lose? (RPO)",
    "Which cloud tools restore service within those limits?",
    "Who declares a disaster — and who executes failover?",
]
for i, q in enumerate(qs):
    top = Inches(1.4 + i * 0.9)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.75), WHITE)
    tb(s, Inches(0.75), top + Inches(0.2), Inches(11.8), Inches(0.4),
       [(f"{i+1}.  {q}", 15, False, INK)])

s = blank()
header(s, "1.4  Cloud Advantages & Limits for DR", "Cloud helps — it does not replace a plan")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.7), "Advantages", [
    "• On-demand capacity in another zone/region",
    "• Managed snapshots and database PITR",
    "• Global load balancers for geo failover",
    "• Pay only for the warm/hot tier you need",
    "• IaC can rebuild stacks quickly",
    "• Easier to run restore drills than tape days",
], TEAL)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.7), "Limits / caveats", [
    "• Cloud ≠ automatic DR for your app data",
    "• Shared responsibility still applies",
    "• Same-region backups can fail together",
    "• Untested restores can miss RTO badly",
    "• Misconfigured IAM can wipe backups too",
    "• Cost rises sharply for near-zero RTO/RPO",
], CORAL)

s = blank()
header(s, "1.5  Shared Responsibility for Disaster Recovery", "Who owns what when systems fail?")
footer(s)
rows = [
    ("Always yours", "App code · data classification · backup policy · IAM · restore testing · runbooks"),
    ("Provider (typical)", "Physical security · underlying hardware · zonal/regional infrastructure reliability"),
    ("Shared design", "You choose multi-zone/region patterns; provider supplies the building blocks"),
    ("Exam reminder", "“We use Google Cloud” is NOT a DR strategy by itself"),
]
for i, (t, d) in enumerate(rows):
    top = Inches(1.35 + i * 1.25)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(1.1), WHITE)
    add_rect(s, Inches(0.5), top, Inches(3.3), Inches(1.1), NAVY)
    tb(s, Inches(0.7), top + Inches(0.35), Inches(2.95), Inches(0.45), [(t, 13, True, WHITE)])
    tb(s, Inches(4.1), top + Inches(0.35), Inches(8.4), Inches(0.5), [(d, 13, False, MUTED)])


# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 2 — Concepts of Backup & DR  (6+ slides)
# ═══════════════════════════════════════════════════════════════════════════

section_divider("02", "Concepts of Backup & Disaster Recovery", [
    "Core definitions",
    "Backup ≠ DR ≠ Business Continuity",
    "Backup types: Full · Incremental · Differential · Snapshot",
    "Backup vs Replication",
    "3-2-1 rule and backup placement",
])

s = blank()
header(s, "2.1  Core Definitions (Part A)", "Terms you must not confuse in exams")
footer(s)
defs = [
    ("Backup", "A copy of data kept so it can be restored after loss, corruption, or deletion."),
    ("Disaster", "An event that makes systems unavailable or data unusable."),
    ("Disaster Recovery", "Policies and procedures to restore IT operations after a disaster."),
    ("Business Continuity", "Keeping critical business services running — broader than IT restore alone."),
]
for i, (t, d) in enumerate(defs):
    top = Inches(1.35 + i * 1.25)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(1.1), WHITE)
    add_rect(s, Inches(0.5), top, Inches(3.5), Inches(1.1), [SKY, CORAL, TEAL, GOLD][i])
    tb(s, Inches(0.7), top + Inches(0.35), Inches(3.15), Inches(0.45), [(t, 13, True, WHITE)])
    tb(s, Inches(4.3), top + Inches(0.3), Inches(8.2), Inches(0.55), [(d, 13, False, MUTED)])

s = blank()
header(s, "2.2  Core Definitions (Part B)", "Restore · Failover · Failback")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(4.05), Inches(4.7), "Restore", [
    "Bringing data/systems back",
    "from a backup or replica",
    "into a usable state.",
    "",
    "Example: create a VM disk",
    "from a GCE snapshot, then",
    "boot and validate the app.",
], SKY)
card(s, Inches(4.7), Inches(1.35), Inches(4.05), Inches(4.7), "Failover", [
    "Switching users from a",
    "failed primary site/system",
    "to a secondary site/system.",
    "",
    "Example: global LB stops",
    "sending traffic to a dead",
    "region and uses the healthy",
    "region instead.",
], TEAL)
card(s, Inches(8.9), Inches(1.35), Inches(4.0), Inches(4.7), "Failback", [
    "Returning operations to the",
    "primary site after it is",
    "repaired and verified.",
    "",
    "Often harder than failover —",
    "data sync and cutover timing",
    "must be careful.",
], CORAL)

s = blank()
header(s, "2.3  Backup ≠ Disaster Recovery", "A zip file is not a recovery plan")
footer(s)
points = [
    ("Only backups", "You still need people, runbooks, spare/rebuildable infra, and a decision of WHEN to fail over."),
    ("Untested restores", "Backups that never restore successfully are false comfort — practice game days."),
    ("Same-region copies", "A region outage can wipe primary AND local backups. Geographic separation matters."),
    ("No RTO/RPO targets", "Without numbers you either overspend on gold-plated DR or underserve a critical app."),
    ("No owner / runbook", "Technology without “who declares disaster” delays every recovery."),
]
for i, (t, d) in enumerate(points):
    top = Inches(1.3 + i * 1.05)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.9), WHITE)
    add_rect(s, Inches(0.5), top, Inches(0.12), Inches(0.9), [CORAL, SKY, TEAL, GOLD, SKY][i])
    tb(s, Inches(0.9), top + Inches(0.12), Inches(11.5), Inches(0.3), [(t, 14, True, NAVY)])
    tb(s, Inches(0.9), top + Inches(0.45), Inches(11.5), Inches(0.35), [(d, 12, False, MUTED)])

s = blank()
header(s, "2.4  Backup Types — Full · Incremental · Differential · Snapshot", "Trade space, time to back up, and restore complexity")
footer(s)
types = [
    ("Full", "Entire dataset every run", "Simple restore\nSlow + expensive\nOften weekly baseline", SKY),
    ("Incremental", "Changes since last backup", "Fast & cheap to take\nRestore = full + chain\nCommon daily schedule", TEAL),
    ("Differential", "Changes since last full", "Restore = full + last diff\nMiddle ground\nKnow for exams", GOLD),
    ("Snapshot", "Point-in-time disk/image", "Near-instant capture\nIdeal for VMs / disks\nGCE snapshots & PVs", CORAL),
]
for i, (t, sub, body, c) in enumerate(types):
    left = Inches(0.35 + i * 3.2)
    add_rounded(s, left, Inches(1.35), Inches(3.05), Inches(4.8), WHITE)
    add_rect(s, left, Inches(1.35), Inches(3.05), Inches(0.65), c)
    tb(s, left + Inches(0.1), Inches(1.45), Inches(2.85), Inches(0.45),
       [(t, 16, True, WHITE)], align=PP_ALIGN.CENTER)
    tb(s, left + Inches(0.15), Inches(2.2), Inches(2.75), Inches(0.5),
       [(sub, 12, True, NAVY)], align=PP_ALIGN.CENTER)
    tb(s, left + Inches(0.2), Inches(3.0), Inches(2.65), Inches(2.5),
       [(ln, 13, False, MUTED) for ln in body.split("\n")], align=PP_ALIGN.CENTER)

s = blank()
header(s, "2.5  Backup vs Replication", "Different tools for different recovery goals")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(3.6), "Backup", [
    "• Point-in-time copies, usually scheduled",
    "• Restore is an explicit action",
    "• Helps roll BACK after corruption / ransomware",
    "• Example: nightly Cloud SQL backup + PITR",
], SKY)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(3.6), "Replication", [
    "• Ongoing copy to another location",
    "• Sync → near-zero RPO, higher cost/latency",
    "• Async → small lag possible, better for distance",
    "• Helps fail OVER to a healthy twin",
], TEAL)
banner(s, Inches(5.25),
       "Replicas can copy corruption. Backups/PITR let you go back in time. Mature designs often use BOTH.")
tb(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.7),
   [("Exam tip: failover to a replica ≠ undo a bad DELETE. Say which problem you are solving.", 13, True, NAVY)])

s = blank()
header(s, "2.6  Where Backups Live & the 3-2-1 Idea", "Placement decides whether backups survive the same disaster")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "Placement levels", [
    "Same zone as primary",
    "  → helps VM failure; NOT zone outage",
    "",
    "Same region, other zone",
    "  → helps zone failure; NOT region outage",
    "",
    "Cross / dual / multi-region",
    "  → survives regional outage better",
    "",
    "Offline / immutable copies",
    "  → stronger against ransomware wipeouts",
], SKY)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "3-2-1 rule (widely taught)", [
    "Keep at least:",
    "• 3 copies of data",
    "• on 2 different media / systems",
    "• with 1 copy off-site",
    "",
    "Cloud mapping:",
    "primary + snapshot/versioning +",
    "another region (or vaulted copy)",
    "",
    "Exact policy varies — exam idea is",
    "redundancy + geographic separation.",
], TEAL)


# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 3 — HA & Fault Tolerance  (6+ slides)
# ═══════════════════════════════════════════════════════════════════════════

section_divider("03", "High Availability & Fault Tolerance", [
    "HA and FT definitions",
    "Uptime percentages",
    "Redundancy & SPOF",
    "HA patterns (instance → region)",
    "HA vs DR — clear exam line",
])

s = blank()
header(s, "3.1  High Availability (HA)", "Keep the service reachable by reducing downtime")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "Definition & goal", [
    "High Availability designs keep the",
    "service reachable for users through",
    "redundancy and automatic failover.",
    "",
    "Goal: minimise downtime users feel.",
    "Usually measured as uptime %.",
], SKY)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "Cloud building blocks", [
    "• Managed Instance Groups (2+ VMs)",
    "• Load balancer + health checks",
    "• Multi-zone deployments",
    "• Kubernetes replicas / Services",
    "• Autohealing of unhealthy members",
    "",
    "HA is proactive design — not luck.",
], TEAL)

s = blank()
header(s, "3.2  Uptime Percentage — What “Nines” Mean", "Often asked conceptually in exams / viva")
footer(s)
# table
headers = ["Uptime", "Approx downtime / year", "Feeling"]
rows = [
    ["99% (2 nines)", "~3.65 days", "Noticeable; many apps reject this"],
    ["99.9% (3 nines)", "~8.8 hours", "Common cloud SLA conversation starter"],
    ["99.99% (4 nines)", "~52 minutes", "Harder; needs strong HA design"],
    ["99.999% (5 nines)", "~5 minutes", "Very expensive / specialised"],
]
for i, h in enumerate(headers):
    left = Inches(0.5 + i * 4.2)
    add_rect(s, left, Inches(1.4), Inches(4.05), Inches(0.55), NAVY)
    tb(s, left + Inches(0.1), Inches(1.5), Inches(3.85), Inches(0.4),
       [(h, 13, True, WHITE)], align=PP_ALIGN.CENTER)
for r, row in enumerate(rows):
    for c, cell in enumerate(row):
        left = Inches(0.5 + c * 4.2)
        top = Inches(2.05 + r * 0.85)
        add_rect(s, left, top, Inches(4.05), Inches(0.8), WHITE if r % 2 == 0 else SOFT)
        tb(s, left + Inches(0.1), top + Inches(0.22), Inches(3.85), Inches(0.4),
           [(cell, 12, c == 0, NAVY if c == 0 else MUTED)], align=PP_ALIGN.CENTER)
tb(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.9),
   [("Note: provider SLA ≠ your application uptime. Your architecture can be worse than the platform SLA.", 13, True, CORAL)])

s = blank()
header(s, "3.3  Fault Tolerance (FT)", "Continue operating correctly when a component fails")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "Definition", [
    "Fault Tolerance means the system",
    "continues to work despite failure of",
    "one or more components — ideally",
    "without the user noticing.",
    "",
    "Examples:",
    "• Autohealing recreates a bad VM",
    "• Replica Pods keep serving traffic",
    "• Stateless apps rebuild from GCS",
], TEAL)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "HA ↔ FT relationship", [
    "• HA often USES fault-tolerant parts",
    "• FT mechanisms help HA targets",
    "• Not every FT design is full regional DR",
    "",
    "Example:",
    "Multi-zone MIG = strong HA + FT",
    "Region gone → still need multi-region",
    "DR + protected data copies",
], SKY)

s = blank()
header(s, "3.4  Redundancy & Single Point of Failure (SPOF)", "If losing ONE thing kills the app — you do not have HA yet")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "Redundancy", [
    "Duplicate components so one loss",
    "does not stop the service.",
    "",
    "Examples:",
    "• 3 Kubernetes Pod replicas",
    "• 2+ VMs in a MIG",
    "• Multi-zone load-balanced backends",
    "• Secondary database replica",
], SKY)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "SPOF examples & fixes", [
    "SPOF → Fix",
    "One VM store → MIG + LB",
    "One zone only → multi-zone",
    "One DB with no backup → PITR + replicas",
    "One DNS entry nowhere else → careful",
    "  multi-region cutover plan",
    "",
    "Find SPOFs before production does.",
], CORAL)

s = blank()
header(s, "3.5  HA Patterns — From One VM to Multi-Region", "Climb this ladder as reliability needs grow")
footer(s)
patterns = [
    ("1. Multi-instance", "MIG / K8s replicas so one crash ≠ outage"),
    ("2. Load balancing", "Send traffic only to healthy backends"),
    ("3. Health checks", "Detect failures early; remove bad nodes"),
    ("4. Autohealing", "Replace unhealthy VMs/Pods automatically"),
    ("5. Multi-zone", "Survive one datacenter failure in a region"),
    ("6. Multi-region", "Survive a whole region outage (true regional DR)"),
]
for i, (t, d) in enumerate(patterns):
    col, row = i % 3, i // 3
    left = Inches(0.45 + col * 4.25)
    top = Inches(1.35 + row * 2.55)
    add_rounded(s, left, top, Inches(4.05), Inches(2.3), WHITE)
    add_rect(s, left, top, Inches(4.05), Inches(0.55), [SKY, TEAL, CORAL, GOLD, SKY, TEAL][i])
    tb(s, left + Inches(0.15), top + Inches(0.12), Inches(3.75), Inches(0.4),
       [(t, 13, True, WHITE)], align=PP_ALIGN.CENTER)
    tb(s, left + Inches(0.2), top + Inches(0.95), Inches(3.65), Inches(1.0),
       [(d, 13, False, MUTED)], align=PP_ALIGN.CENTER)

s = blank()
header(s, "3.6  HA vs DR — Clear Line for Exams", "Both needed; they solve different classes of failure")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.0), "HA focus", [
    "Keep the app up during common failures:",
    "VM death, process crash, single zone blip.",
    "",
    "Typical tools: MIG, LB, health checks,",
    "replicas, multi-zone.",
], SKY)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.0), "DR focus", [
    "Recover after major loss:",
    "region gone, ransomware, data corruption.",
    "",
    "Typical tools: cross-region replicas,",
    "backups/PITR, runbooks, strategy tiers.",
], CORAL)
banner(s, Inches(5.7),
       "Food delivery apps need HA every dinner hour. They STILL need DR when a region or database event hits.")


# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 4 — DR Strategies RPO & RTO  (10+ slides)
# ═══════════════════════════════════════════════════════════════════════════

section_divider("04", "Disaster Recovery Strategies (RPO & RTO)", [
    "RTO and RPO definitions",
    "Timeline view",
    "How metrics drive cost",
    "Strategy ladder cold → hot",
    "Deep dive each strategy + business choice",
])

s = blank()
header(s, "4.1  RTO — Recovery Time Objective", "Maximum acceptable downtime")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "Definition", [
    "RTO = max time the business accepts",
    "being offline after a disaster.",
    "",
    "Question: “How long can we be down?”",
    "",
    "Measured from disaster → restored",
    "service (users can work again).",
], CORAL)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "Examples", [
    "UPI / payments: seconds to few minutes",
    "Hospital EHR: very low RTO (life safety)",
    "Campus result portal: may accept hours",
    "  except on publish day (then tighten)",
    "Weekly internal report: hours–a day OK",
    "",
    "RTO is a BUSINESS decision first,",
    "then an engineering design constraint.",
], SKY)

s = blank()
header(s, "4.2  RPO — Recovery Point Objective", "Maximum acceptable data loss (as a time window)")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "Definition", [
    "RPO = max age of data you can afford",
    "to lose when you restore.",
    "",
    "Question: “How much recent work can",
    "we lose?”",
    "",
    "Measured as time looking BACK from",
    "the disaster to the last good copy.",
], SKY)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "Examples", [
    "Hourly backups ⇒ RPO up to ~1 hour",
    "Continuous async replication ⇒ minutes",
    "Synchronous replication ⇒ near zero",
    "Cloud SQL PITR ⇒ often minutes-level",
    "",
    "Tighter RPO needs more frequent",
    "protection — and usually more cost.",
], TEAL)

s = blank()
header(s, "4.3  RPO & RTO on a Timeline", "Never swap these two in exam answers")
footer(s)
add_rounded(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(3.2), WHITE)
add_rect(s, Inches(1.2), Inches(3.1), Inches(10.0), Inches(0.08), SOFT)
events = [
    (1.2, "Last good\nbackup / replica", TEAL),
    (5.0, "DISASTER", CORAL),
    (8.5, "Service\nrestored", SKY),
    (11.2, "Users\nback", GREEN),
]
for x, label, c in events:
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.9), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = c
    circ.line.fill.background()
    tb(s, Inches(x - 0.55), Inches(3.55), Inches(1.7), Inches(0.8),
       [(label, 11, True, NAVY)], align=PP_ALIGN.CENTER)
tb(s, Inches(1.0), Inches(1.4), Inches(11), Inches(0.35),
   [("RPO = last good copy → disaster   (data you may lose)", 14, True, TEAL)])
tb(s, Inches(1.0), Inches(5.3), Inches(11), Inches(0.35),
   [("RTO = disaster → service restored   (downtime users feel)", 14, True, CORAL)])
tb(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.7),
   [("Memory trick: RPO looks into the Past (data). RTO looks toward Recovery Time (clock).", 13, True, NAVY)])

s = blank()
header(s, "4.4  How RTO/RPO Drive Cost", "Lower numbers → climb the strategy ladder → spend more")
footer(s)
simple_lines = [
    "Lower RTO → must failover/rebuild faster → more automation + more standby capacity",
    "Lower RPO → protect data more often → frequent backups / continuous replication",
    "Near-zero RTO + near-zero RPO ≈ Hot Standby or Active-Active (expensive)",
    "Loose RTO/RPO ≈ Backup & Restore may be enough (cheapest)",
    "Choose the cheapest strategy that STILL meets the business numbers",
    "Do not buy banking-grade DR for a weekly PDF report job",
]
for i, line in enumerate(simple_lines):
    top = Inches(1.35 + i * 0.8)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.65), WHITE)
    add_rect(s, Inches(0.5), top, Inches(0.12), Inches(0.65), [SKY, TEAL, CORAL, GOLD, SKY, TEAL][i])
    tb(s, Inches(0.9), top + Inches(0.15), Inches(11.6), Inches(0.4), [(line, 13, False, INK)])

s = blank()
header(s, "4.5  DR Strategy Ladder (Cold → Hot)", "Five tiers every Unit II student must memorise")
footer(s)
headers2 = ["Strategy", "Typical RTO", "Typical RPO", "Cost feel"]
rows2 = [
    ["Backup & Restore", "Hours", "Hours", "Lowest"],
    ["Pilot Light", "Tens of minutes", "Minutes", "Low–medium"],
    ["Warm Standby", "Minutes", "Minutes", "Medium"],
    ["Hot Standby", "Near zero", "Near zero", "High"],
    ["Active-Active", "~Zero", "~Zero", "Highest"],
]
for i, h in enumerate(headers2):
    left = Inches(0.4 + i * 3.2)
    add_rect(s, left, Inches(1.35), Inches(3.05), Inches(0.5), NAVY)
    tb(s, left + Inches(0.05), Inches(1.42), Inches(2.95), Inches(0.35),
       [(h, 12, True, WHITE)], align=PP_ALIGN.CENTER)
for r, row in enumerate(rows2):
    for c, cell in enumerate(row):
        left = Inches(0.4 + c * 3.2)
        top = Inches(1.95 + r * 0.8)
        add_rect(s, left, top, Inches(3.05), Inches(0.75), WHITE if r % 2 == 0 else SOFT)
        color = [SKY, TEAL, CORAL, GOLD, GREEN][r] if c == 0 else INK
        tb(s, left + Inches(0.05), top + Inches(0.2), Inches(2.95), Inches(0.4),
           [(cell, 12, c == 0, color)], align=PP_ALIGN.CENTER)
tb(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
   [("Next slides = deep dive into each tier.", 13, True, MUTED)])

s = blank()
header(s, "4.6  Strategy Deep Dive — Backup & Restore", "Coldest / cheapest tier")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "How it works", [
    "1. Take periodic backups (disk/DB/object)",
    "2. On disaster: provision infrastructure",
    "3. Restore data from backup",
    "4. Validate application",
    "5. Point users (DNS / LB) to recovered stack",
    "",
    "RTO/RPO: often hours",
    "Cost: lowest ongoing",
], SKY)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "Fit & risks", [
    "Good for:",
    "• Blogs, labs, internal tools",
    "• Non-critical reporting",
    "• Systems that can wait",
    "",
    "Risks:",
    "• Long downtime if restore untested",
    "• Easy to believe backups exist when",
    "  they do not restore cleanly",
], CORAL)

s = blank()
header(s, "4.7  Strategy Deep Dive — Pilot Light", "Tiny core always on; scale out when disaster hits")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "How it works", [
    "• Keep critical data tier alive in DR site",
    "  (e.g. DB replica + minimal services)",
    "• Most compute stays off until needed",
    "• On disaster: scale fleets + cut over",
    "",
    "RTO: tens of minutes (practiced)",
    "RPO: minutes (depends on replication)",
], TEAL)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "Fit & risks", [
    "Good for:",
    "• Business apps that can wait tens of",
    "  minutes but not hours",
    "",
    "Needs:",
    "• Clear scale-out automation",
    "• Runbook skill / game days",
    "• Monitoring of the “pilot” core",
], SKY)

s = blank()
header(s, "4.8  Strategy Deep Dive — Warm Standby", "Scaled-down full copy always running")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "How it works", [
    "• Complete environment runs smaller",
    "  in the DR site continuously",
    "• Data kept fairly fresh",
    "• Failover ≈ traffic shift + scale up",
    "",
    "RTO/RPO: typically minutes",
    "Cost: medium (always-on smaller stack)",
], SKY)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "Fit", [
    "Good for:",
    "• Medium-critical business apps",
    "• Campus ERP during exam windows",
    "• Systems where hours are unacceptable",
    "  but full hot duplicate is too costly",
    "",
    "Must still test failover — warm is not",
    "automatic excellence.",
], TEAL)

s = blank()
header(s, "4.9  Strategy Deep Dive — Hot Standby", "Near full-size duplicate ready to take over")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "How it works", [
    "• Secondary almost as capable as primary",
    "• Failover measured in seconds–low minutes",
    "• Continuous (often sync/async) data protect",
    "",
    "RTO/RPO: near zero",
    "Cost: high ≈ running two productions",
], CORAL)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "Fit", [
    "Good for:",
    "• Payments / trading-style systems",
    "• Life-critical records access",
    "• Brand-critical customer checkouts",
    "",
    "Watch:",
    "• Data consistency during failover",
    "• Failback planning after primary returns",
], SKY)

s = blank()
header(s, "4.10  Strategy Deep Dive — Multi-site Active-Active", "Multiple regions serve live traffic together")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "How it works", [
    "• Two+ regions serve users at once",
    "• Global HTTP(S) LB + often CDN",
    "• One region dies → others continue",
    "",
    "RTO/RPO: ~zero when designed well",
    "Hardest: data consistency & conflicts",
], TEAL)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "Fit & warnings", [
    "Good for:",
    "• Global SaaS / major commerce",
    "• Ultra-low downtime budgets",
    "",
    "Warnings:",
    "• Databases across regions are hard",
    "• Not every app is active-active ready",
    "• Highest cost and design complexity",
], CORAL)

s = blank()
header(s, "4.11  Choosing a Strategy — Business First", "Same cloud tools · different RTO/RPO budgets")
footer(s)
choices = [
    ("Campus ERP / results", "Warm/Hot in result week; colder in off-season"),
    ("Food delivery API", "Hot / Active-Active mindset — minutes kill orders"),
    ("Hospital patient records", "Hot + strict RPO; audited restores; multi-region thinking"),
    ("Batch analytics / ML train", "Backup & Restore + retries/checkpoints often enough"),
    ("UPI / payments", "Very low RTO/RPO; multi-zone HA + strong DR drills"),
]
for i, (t, d) in enumerate(choices):
    top = Inches(1.3 + i * 1.0)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.85), WHITE)
    add_rect(s, Inches(0.5), top, Inches(4.0), Inches(0.85), NAVY)
    tb(s, Inches(0.7), top + Inches(0.25), Inches(3.6), Inches(0.4), [(t, 13, True, WHITE)])
    tb(s, Inches(4.8), top + Inches(0.25), Inches(7.7), Inches(0.45), [(d, 13, False, MUTED)])


# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 5 — Cloud Backups & Snapshots  (7+ slides)
# ═══════════════════════════════════════════════════════════════════════════

section_divider("05", "Cloud-Based Backup Solutions & Snapshots", [
    "GCP building blocks overview",
    "Compute Engine snapshots",
    "Recovery playbook",
    "Cloud Storage versioning",
    "Cloud SQL backups & PITR",
    "GKE / PaaS data notes + best practices",
])

s = blank()
header(s, "5.1  Cloud Backup Building Blocks (GCP View)", "Same ideas exist on AWS/Azure under different names")
footer(s)
sols = [
    ("Compute Engine\nSnapshots", "Disk point-in-time\nIncremental after first\nRestore → new disk/VM", SKY),
    ("Cloud Storage", "Object versioning\nDual/multi-region\nLifecycle to cold tiers", TEAL),
    ("Cloud SQL", "Automated backups\nPoint-in-time recovery\nTight DB RPO", CORAL),
    ("GKE / Apps", "PV snapshots\nMulti-cluster failover\nApp-level backups still yours", GOLD),
]
for i, (t, body, c) in enumerate(sols):
    left = Inches(0.35 + i * 3.2)
    add_rounded(s, left, Inches(1.35), Inches(3.05), Inches(4.8), WHITE)
    add_rect(s, left, Inches(1.35), Inches(3.05), Inches(1.05), c)
    tb(s, left + Inches(0.1), Inches(1.45), Inches(2.85), Inches(0.9),
       [(t, 13, True, WHITE)], align=PP_ALIGN.CENTER)
    tb(s, left + Inches(0.15), Inches(2.7), Inches(2.75), Inches(2.8),
       [(ln, 13, False, MUTED) for ln in body.split("\n")], align=PP_ALIGN.CENTER)

s = blank()
header(s, "5.2  Compute Engine Snapshots — What They Are", "Point-in-time disk copies for VM recovery")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "Properties", [
    "• Capture a persistent disk at a moment",
    "• First snapshot = full baseline",
    "• Later snapshots = incremental",
    "• Stored durably (Cloud Storage backed)",
    "• Create new disk / VM from a snapshot",
    "",
    "Great for VM/boot disk recovery patterns.",
], SKY)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "Command shape", [
    "gcloud compute disks snapshot DISK \\",
    "  --snapshot-names=NAME \\",
    "  --zone=ZONE",
    "",
    "Also automate via schedules / policies",
    "in production — do not rely on memory.",
    "",
    "Caution: busy DB disks may be only",
    "crash-consistent — prefer managed DB",
    "backups for transactional engines.",
], TEAL)

s = blank()
header(s, "5.3  Snapshot Recovery Playbook", "Practice this BEFORE the outage")
footer(s)
steps = [
    ("1", "Identify", "Pick last known-good snapshot"),
    ("2", "Create disk", "New disk from snapshot"),
    ("3", "Attach / boot", "Create VM or attach disk"),
    ("4", "Validate", "App checks + data sanity"),
    ("5", "Cut over", "Point LB / DNS if needed"),
    ("6", "Measure", "Compare elapsed time vs RTO"),
]
for i, (n, t, d) in enumerate(steps):
    col, row = i % 3, i // 3
    left = Inches(0.45 + col * 4.25)
    top = Inches(1.35 + row * 2.55)
    add_rounded(s, left, top, Inches(4.05), Inches(2.3), WHITE)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.6), top + Inches(0.25),
                              Inches(0.7), Inches(0.7))
    circ.fill.solid()
    circ.fill.fore_color.rgb = [SKY, TEAL, CORAL, GOLD, SKY, TEAL][i]
    circ.line.fill.background()
    tb(s, left + Inches(1.6), top + Inches(0.35), Inches(0.7), Inches(0.5),
       [(n, 16, True, WHITE)], align=PP_ALIGN.CENTER)
    tb(s, left + Inches(0.2), top + Inches(1.15), Inches(3.65), Inches(0.35),
       [(t, 14, True, NAVY)], align=PP_ALIGN.CENTER)
    tb(s, left + Inches(0.2), top + Inches(1.55), Inches(3.65), Inches(0.5),
       [(d, 12, False, MUTED)], align=PP_ALIGN.CENTER)

s = blank()
header(s, "5.4  Cloud Storage — Versioning & Regional Resilience", "Protect objects, dumps, media, and artifacts")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "Key features", [
    "• Object versioning → undelete / prior versions",
    "• Dual-region / multi-region buckets",
    "• Lifecycle policies → Nearline / Coldline / Archive",
    "• Soft delete / retention options (know conceptually)",
    "",
    "Use for: DB dumps, media, ML datasets,",
    "build artifacts, log archives.",
], SKY)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "DR notes", [
    "• Single-region bucket can share fate",
    "  with an application region outage",
    "• Dual/multi-region improves survival of",
    "  object data across regional events",
    "• Versioning helps accidental overwrite",
    "• Restrict IAM — backups are sensitive",
], CORAL)

s = blank()
header(s, "5.5  Cloud SQL Backups & Point-in-Time Recovery", "Transactional data needs service-native protection")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "Capabilities to know", [
    "• Automated scheduled backups",
    "• On-demand backup before risky changes",
    "• Point-in-Time Recovery (PITR)",
    "  → restore to a timestamp",
    "• Export options to Cloud Storage",
], SKY)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "Classic exam story", [
    "Corruption / bad DELETE at 2:00 PM",
    "Business needs ≈ 1:30 PM state",
    "",
    "PITR restores to that timestamp.",
    "",
    "Replicas alone may have copied the",
    "bad state — PITR rolls BACK in time.",
], TEAL)
banner(s, Inches(6.2),
       "Prefer managed DB backups/PITR over raw VM disk snapshots for databases whenever available.")

s = blank()
header(s, "5.6  GKE & PaaS — What You Still Own", "Managed platforms ≠ managed application data strategy")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "GKE", [
    "• Control plane etcd managed by Google",
    "• Your app state on Persistent Volumes",
    "  still needs snapshot/backup policy",
    "• Multi-cluster / multi-region helps",
    "  compute failover — pair with data DR",
], TEAL)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "App Engine / Functions", [
    "• Platform HA is largely Google-managed",
    "• YOUR data in Cloud SQL / Firestore /",
    "  GCS still needs backup & DR policy",
    "• PaaS does not erase shared",
    "  responsibility for application data",
], SKY)

s = blank()
header(s, "5.7  Cloud Backup Best Practices", "Production checklist (and exam favourites)")
footer(s)
practices = [
    "Automate schedules — do not rely on human memory",
    "Alert when backups fail; treat silence as danger",
    "Encrypt backups; restrict IAM with least privilege",
    "Keep at least one copy outside the primary region",
    "Test restores on a calendar; record actual RTO",
    "Use versioning / retention that survives mistakes",
    "For ransomware awareness: consider immutable / vaulted copies",
    "Document restore owners and communication steps",
]
for i, p in enumerate(practices):
    top = Inches(1.3 + i * 0.65)
    add_rect(s, Inches(0.5), top + Inches(0.08), Inches(0.16), Inches(0.16),
             [TEAL, SKY][i % 2])
    tb(s, Inches(0.9), top, Inches(11.8), Inches(0.5), [(p, 13, False, INK)])


# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 6 — Case Studies  (9+ slides)
# ═══════════════════════════════════════════════════════════════════════════

section_divider("06", "Case Studies on Disaster Recovery Implementation", [
    "Lab pattern: Fancy Store HA",
    "Global e-commerce regional outage",
    "Data-critical PITR",
    "Batch analytics fault tolerance",
    "Industry: streaming · payments · GitLab · campus ERP",
])

s = blank()
header(s, "6.0  Case Study Map", "Use these stories in long answers / viva")
footer(s)
cases = [
    ("01", "Fancy Store (Lab)", "SPOF VM → MIG + LB + heal"),
    ("02", "Global E-commerce", "Multi-region + global LB + CDN"),
    ("03", "Data-Critical App", "Cloud SQL PITR for minutes RPO"),
    ("04", "Batch Analytics", "Spot + retries + GCS state"),
    ("05", "Hotstar / Streaming", "Demand spike resilience"),
    ("06", "UPI / Payments", "Strict RTO/RPO thinking"),
    ("07", "GitLab 2017", "Backup gaps & restore lessons"),
    ("08", "Campus ERP", "Time-varying DR tier"),
]
for i, (n, t, d) in enumerate(cases):
    col, row = i % 4, i // 4
    left = Inches(0.35 + col * 3.2)
    top = Inches(1.35 + row * 2.55)
    add_rounded(s, left, top, Inches(3.05), Inches(2.3), WHITE)
    add_rect(s, left, top, Inches(3.05), Inches(0.55), [SKY, TEAL, CORAL, GOLD, SKY, TEAL, CORAL, GOLD][i])
    tb(s, left + Inches(0.1), top + Inches(0.12), Inches(2.85), Inches(0.4),
       [(n + "  " + t, 12, True, WHITE)], align=PP_ALIGN.CENTER)
    tb(s, left + Inches(0.15), top + Inches(0.9), Inches(2.75), Inches(1.1),
       [(d, 12, False, MUTED)], align=PP_ALIGN.CENTER)

s = blank()
header(s, "6.1  Case — Fancy Store on Compute Engine (Lab)", "Problem → Solution → Level achieved")
footer(s)
card(s, Inches(0.4), Inches(1.35), Inches(4.1), Inches(4.7), "Problem", [
    "Single VM hosts the store.",
    "VM dies → whole store down.",
    "",
    "Classic SPOF.",
    "No health checks,",
    "no automatic recovery.",
], CORAL)
card(s, Inches(4.65), Inches(1.35), Inches(4.1), Inches(4.7), "Solution", [
    "• Managed Instance Group",
    "• Health checks",
    "• Autohealing",
    "• HTTP(S) Load Balancer",
    "• Named ports for services",
    "",
    "Identical VMs from templates",
    "(cattle, not pets).",
], SKY)
card(s, Inches(8.9), Inches(1.35), Inches(4.0), Inches(4.7), "Level", [
    "Strong HA in a zone.",
    "",
    "Better with multi-zone MIG.",
    "",
    "Still need data backups",
    "for true DR events",
    "(corruption / region loss).",
], TEAL)

s = blank()
header(s, "6.2  Case — Global E-commerce Regional Outage", "When single-region HA is not enough")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(3.8), "Problem", [
    "All capacity in one region.",
    "Regional outage stops checkout for everyone.",
    "Revenue and brand damage are immediate.",
], CORAL)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(3.8), "Solution", [
    "• Multi-region deployment",
    "• Global HTTP(S) Load Balancer",
    "• Cloud CDN for static assets",
    "• Warm/Hot or Active-Active + data replication",
], SKY)
banner(s, Inches(5.5),
       "Amazon / Flipkart sale peaks and Shopify merchants plan regional failover so one region wobble doesn’t cancel payments.")

s = blank()
header(s, "6.3  Case — Data-Critical Application (PITR)", "Corruption at 2:00 PM → restore to 1:30 PM")
footer(s)
steps = [
    ("1", "Incident", "Bad write / DELETE /\nransomware-like event"),
    ("2", "RPO decision", "Business: lose ≤ 30 min\nof transactions"),
    ("3", "PITR restore", "Cloud SQL restore to\n1:30 PM state"),
    ("4", "Validate", "App checks, reconcile,\nreopen traffic"),
]
for i, (n, t, d) in enumerate(steps):
    left = Inches(0.4 + i * 3.2)
    add_rounded(s, left, Inches(1.4), Inches(3.05), Inches(3.4), WHITE)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.15), Inches(1.65),
                              Inches(0.7), Inches(0.7))
    circ.fill.solid()
    circ.fill.fore_color.rgb = [CORAL, SKY, TEAL, GOLD][i]
    circ.line.fill.background()
    tb(s, left + Inches(1.15), Inches(1.75), Inches(0.7), Inches(0.5),
       [(n, 16, True, WHITE)], align=PP_ALIGN.CENTER)
    tb(s, left + Inches(0.15), Inches(2.55), Inches(2.75), Inches(0.4),
       [(t, 14, True, NAVY)], align=PP_ALIGN.CENTER)
    tb(s, left + Inches(0.15), Inches(3.15), Inches(2.75), Inches(1.2),
       [(d, 12, False, MUTED)], align=PP_ALIGN.CENTER)
tb(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.2),
   [("Teaching point: failover to a replica can copy the poison. PITR rolls BACK in time — different problem.", 14, True, NAVY),
    ("Maps to: college marksheets, inventory DBs, fintech ledgers.", 13, False, MUTED)])

s = blank()
header(s, "6.4  Case — Batch Analytics / Non-Interactive Jobs", "Not every workload needs Hot Standby")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "Problem & solution", [
    "Problem: job fails when a Spot node dies",
    "",
    "Solution:",
    "• Keep checkpoints / state in Cloud Storage",
    "• Retry / resume from last watermark",
    "• Accept longer RTO if business allows",
    "",
    "DR level: fault-tolerant & cost-aware",
], TEAL)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "Real parallels", [
    "• ML training with periodic checkpoints",
    "• Nightly ETL pipelines",
    "• Log processing with replay",
    "• Recommendation batch builds",
    "",
    "Design for interruption — not immortality.",
    "Save money where humans are not waiting.",
], SKY)

s = blank()
header(s, "6.5  Case — Streaming Spike (Hotstar / Live Sports Style)", "Disaster of demand is still an availability problem")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.0), "Problem", [
    "Traffic multiplies hugely during live events.",
    "Failure mode may be overload — not only",
    "datacenter loss.",
], CORAL)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.0), "Solution themes", [
    "• Autoscaling compute",
    "• CDN near viewers",
    "• Multi-AZ resilient origins",
    "• Graceful degradation (lower quality",
    "  beats total blackout)",
], SKY)
banner(s, Inches(5.7),
       "Goal during the last over: stay up. Backups alone do not stream a match.")

s = blank()
header(s, "6.6  Case — Payments / UPI Style Systems", "When minutes of downtime destroy trust")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "Requirements", [
    "• Very low RTO (transactions must flow)",
    "• Very low RPO (money must not vanish)",
    "• Multi-zone HA as the baseline",
    "• Strict failover drills / game days",
    "• Careful DB replication + backup strategy",
], CORAL)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "Teaching point", [
    "Not every campus app needs banking-grade",
    "Active-Active spend.",
    "",
    "But you must explain WHY payments pay",
    "for hotter tiers — and map that to RTO/RPO.",
    "",
    "Salary-day spikes also need capacity,",
    "not only recovery docs.",
], SKY)

s = blank()
header(s, "6.7  Case — GitLab Database Incident (2017)", "Industry lesson: backups that cannot restore are theatre")
footer(s)
points = [
    ("What happened (high level)", "Production database deletion event + backup/restore difficulties caused a long public outage."),
    ("Lesson 1", "Multiple backup methods beat a single “we think backups exist” assumption."),
    ("Lesson 2", "Restore testing matters as much as backup scheduling."),
    ("Lesson 3", "Clear RPO expectations and access controls reduce panic and blast radius."),
]
for i, (t, d) in enumerate(points):
    top = Inches(1.35 + i * 1.25)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(1.1), WHITE)
    add_rect(s, Inches(0.5), top, Inches(0.12), Inches(1.1), [CORAL, SKY, TEAL, GOLD][i])
    tb(s, Inches(0.9), top + Inches(0.18), Inches(11.5), Inches(0.3), [(t, 14, True, NAVY)])
    tb(s, Inches(0.9), top + Inches(0.55), Inches(11.5), Inches(0.4), [(d, 13, False, MUTED)])

s = blank()
header(s, "6.8  Case — Campus ERP / Exam Result Portal", "RTO/RPO can be time-varying")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.05), Inches(4.6), "Problem", [
    "Most of year: moderate traffic.",
    "Result day: load explodes; any outage",
    "causes campus-wide chaos.",
    "",
    "Database integrity also critical before",
    "publish.",
], CORAL)
card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(4.6), "Practical approach", [
    "Off-peak: Backup & Restore + basic HA",
    "Peak window:",
    "• Raise capacity / warmer standby",
    "• CDN for static result pages",
    "• Queues / rate limits",
    "• Fresh tested DB backup before publish",
    "",
    "Escalate DR tier for critical windows.",
], TEAL)


# ═══════════════════════════════════════════════════════════════════════════
# CLOSING
# ═══════════════════════════════════════════════════════════════════════════

s = blank()
header(s, "People & Process — Often Forgotten Subtopic", "Technology without ownership delays every recovery")
footer(s)
qs = [
    "Who is authorised to declare a disaster?",
    "Who executes failover vs who informs users / management?",
    "Where is the runbook, and is it reachable if Slack/primary tools are down?",
    "How do we fail back safely after primary recovers?",
    "When was the last successful restore test — and what RTO did we actually hit?",
    "Are secrets / IAM for backups separate enough to survive ransomware credentials?",
]
for i, q in enumerate(qs):
    top = Inches(1.3 + i * 0.85)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.7), WHITE)
    tb(s, Inches(0.75), top + Inches(0.18), Inches(11.8), Inches(0.4),
       [(f"{i+1}.  {q}", 13, False, INK)])

s = blank()
header(s, "Exam Rapid Revision — Sheet A", "Definitions & distinctions")
footer(s)
lines = [
    "Backup = data copy · Disaster = unavailability/data-loss event · DR = IT restore plan · Continuity = business keeps running",
    "Restore / Failover / Failback are different actions — learn all three",
    "Backup types: Full, Incremental, Differential, Snapshot",
    "Backup ≠ Replication; replicas can copy corruption; PITR rolls back in time",
    "HA minimises downtime; FT continues despite failure; kill SPOFs with redundancy",
    "RTO = downtime budget · RPO = data-loss budget — never swap",
]
for i, line in enumerate(lines):
    top = Inches(1.35 + i * 0.85)
    add_rect(s, Inches(0.5), top + Inches(0.1), Inches(0.16), Inches(0.16), SKY if i % 2 == 0 else TEAL)
    tb(s, Inches(0.9), top, Inches(11.8), Inches(0.7), [(line, 13, False, INK)])

s = blank()
header(s, "Exam Rapid Revision — Sheet B", "Strategies · tools · cases")
footer(s)
lines = [
    "Ladder: Backup & Restore → Pilot Light → Warm → Hot → Active-Active",
    "Lower RTO/RPO ⇒ higher cost; choose cheapest that still meets business numbers",
    "GCP: GCE snapshots · GCS versioning/dual-multi-region · Cloud SQL backups+PITR · GKE PV / multi-cluster",
    "HA ≠ full DR: multi-zone MIG still needs cross-region data protection for regional loss",
    "Cases: Fancy Store HA · Global multi-region · SQL PITR · Batch FT · Streaming · UPI · GitLab · Campus ERP",
    "Always: automate backups, alert failures, test restores, document owners",
]
for i, line in enumerate(lines):
    top = Inches(1.35 + i * 0.85)
    add_rect(s, Inches(0.5), top + Inches(0.1), Inches(0.16), Inches(0.16), CORAL if i % 2 == 0 else GOLD)
    tb(s, Inches(0.9), top, Inches(11.8), Inches(0.7), [(line, 13, False, INK)])

s = blank()
header(s, "Key Takeaways", "Design DR for people and money — not for slides")
footer(s)
takes = [
    ("Measure first", "Agree RTO & RPO with stakeholders, then pick the cheapest strategy that meets them."),
    ("HA ≠ full DR", "Multi-zone MIG is HA. Surviving regional loss needs multi-region design + protected data."),
    ("Practice restores", "Untested backups are theatre. Game days prove RTO; PITR drills prove RPO."),
    ("Match the workload", "UPI ≠ nightly ETL. Active-Active for payments; retries + GCS for batch."),
    ("Own the process", "Declare ownership, runbooks, failback, and backup IAM — technology alone is incomplete."),
]
for i, (t, d) in enumerate(takes):
    top = Inches(1.3 + i * 1.05)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.9), WHITE)
    add_rect(s, Inches(0.5), top, Inches(0.12), Inches(0.9), [SKY, TEAL, CORAL, GOLD, SKY][i])
    tb(s, Inches(0.9), top + Inches(0.12), Inches(11.5), Inches(0.3), [(t, 14, True, NAVY)])
    tb(s, Inches(0.9), top + Inches(0.45), Inches(11.5), Inches(0.35), [(d, 12, False, MUTED)])

s = blank()
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.15), TEAL)
tb(s, Inches(0.8), Inches(2.1), Inches(11.5), Inches(0.8), [("Thank you", 44, True, WHITE)])
tb(s, Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.5),
   [("Questions · Discussion · Map each case to a DR tier", 18, False, SKY_LIGHT)])
tb(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(1.4),
   [("Unit II — Disaster Recovery in Cloud (expanded subtopic deck)", 14, False, SOFT),
    ("Pair with: CC-II-Unit-II-Disaster-Recovery-in-Cloud.txt", 14, False, SOFT),
    ("FIT & CS, Parul University", 14, False, MUTED)])
PAGE += 1


out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "CC-II-Unit-II-Disaster-Recovery-in-Cloud.pptx")
prs.save(out)
print(f"Saved {out}")
print(f"Slides: {len(prs.slides)}")
