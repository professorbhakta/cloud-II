#!/usr/bin/env python3
"""
Unit IV PowerPoint — Kubernetes & Container Orchestration
Expanded deck: 4–5+ slides per syllabus topic with subtopics + real-life cases.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

NAVY = RGBColor(0x0A, 0x1B, 0x2E)
SKY = RGBColor(0x1A, 0x6F, 0xB5)
SKY_LIGHT = RGBColor(0x3D, 0x9B, 0xE0)
TEAL = RGBColor(0x0D, 0xA5, 0xA0)
CORAL = RGBColor(0xE8, 0x6A, 0x4C)
CLOUD = RGBColor(0xF4, 0xF8, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xE8, 0xF1, 0xF8)
INK = RGBColor(0x1C, 0x2B, 0x3A)
MUTED = RGBColor(0x5A, 0x6B, 0x7C)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
GREEN = RGBColor(0x1F, 0x8A, 0x5F)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
PAGE = 0


def add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_rounded(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def set_run(run, text, size=18, bold=False, color=INK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def tb(slide, left, top, width, height, lines, default_size=16, default_color=INK,
       align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, bold, color = item, default_size, False, default_color
        else:
            text = item[0]
            size = item[1] if len(item) > 1 else default_size
            bold = item[2] if len(item) > 2 else False
            color = item[3] if len(item) > 3 else default_color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(3)
        run = p.add_run()
        set_run(run, text, size=size, bold=bold, color=color)
    return box


def footer(slide):
    global PAGE
    PAGE += 1
    add_rect(slide, Inches(0), Inches(7.15), SLIDE_W, Inches(0.35), NAVY)
    tb(slide, Inches(0.35), Inches(7.18), Inches(10.5), Inches(0.3),
       [("CC-II  ·  Unit IV  ·  Kubernetes & Container Orchestration  ·  Parul University",
         10, False, SOFT)])
    tb(slide, Inches(11.6), Inches(7.18), Inches(1.4), Inches(0.3),
       [(str(PAGE), 10, False, SOFT)], align=PP_ALIGN.RIGHT)


def header(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.08), TEAL)
    tb(slide, Inches(0.45), Inches(0.2), Inches(12.3), Inches(0.45),
       [(title, 24, True, WHITE)])
    if subtitle:
        tb(slide, Inches(0.45), Inches(0.68), Inches(12.3), Inches(0.32),
           [(subtitle, 12, False, SKY_LIGHT)])


def blank():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CLOUD)
    return slide


def card(slide, left, top, width, height, title, body_lines, accent=SKY):
    add_rounded(slide, left, top, width, height, WHITE)
    add_rect(slide, left, top, Inches(0.1), height, accent)
    tb(slide, left + Inches(0.22), top + Inches(0.12), width - Inches(0.35), Inches(0.35),
       [(title, 14, True, NAVY)])
    formatted = [(ln, 12, False, MUTED) if isinstance(ln, str) else ln for ln in body_lines]
    tb(slide, left + Inches(0.22), top + Inches(0.5), width - Inches(0.35), height - Inches(0.6),
       formatted)


def banner(slide, top, text, color=CORAL):
    add_rounded(slide, Inches(0.45), top, Inches(12.4), Inches(0.5), RGBColor(0xFF, 0xF4, 0xEC))
    tb(slide, Inches(0.65), top + Inches(0.08), Inches(12), Inches(0.35),
       [(f"Real life:  {text}", 12, True, color)])


def section_divider(section_no, title, bullets):
    s = blank()
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, TEAL)
    tb(s, Inches(0.8), Inches(2.0), Inches(11), Inches(0.4),
       [(f"TOPIC {section_no}", 14, True, TEAL)])
    tb(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.0),
       [(title, 34, True, WHITE)])
    for i, b in enumerate(bullets):
        tb(s, Inches(0.8), Inches(4.0 + i * 0.4), Inches(11.5), Inches(0.35),
           [(f"→  {b}", 14, False, SOFT)])
    footer(s)
    return s


# ═══════════════════════════════════════════════════════════════════════════
# OPENING
# ═══════════════════════════════════════════════════════════════════════════

s = blank()
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.15), TEAL)
tb(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.4),
   [("CLOUD COMPUTING – II  ·  UNIT IV", 15, True, TEAL)])
tb(s, Inches(0.8), Inches(2.1), Inches(11.5), Inches(1.3),
   [("Kubernetes & Container", 40, True, WHITE),
    ("Orchestration", 40, True, WHITE)])
tb(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(1.2),
   [("Complete classroom deck  ·  Syllabus topics with subtopics", 14, False, SOFT),
    ("Architecture  ·  Cluster setup  ·  Deploy  ·  Networking  ·  Scale", 14, False, SOFT),
    ("GKE · kubectl · Services · Ingress · HPA · exam revision", 14, False, SOFT)])
tb(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.4),
   [("FIT & CS  ·  Parul University  ·  Weightage 15%  ·  7 hours", 14, False, MUTED)])
PAGE += 1

s = blank()
header(s, "Learning Outcomes", "By the end of Unit IV you should be able to…")
footer(s)
outcomes = [
    ("Define Kubernetes", "Explain why containers need an orchestrator at cluster scale."),
    ("Draw architecture", "Control plane vs worker nodes — API Server, Scheduler, etcd, kubelet."),
    ("Set up a cluster", "Create & connect to GKE (Autopilot/Standard); use namespaces."),
    ("Deploy apps", "Image → registry → Deployment → Service; rolling update & rollback."),
    ("Expose traffic", "Service types + Ingress; why Pod IPs are ephemeral."),
    ("Monitor & scale", "HPA, cluster capacity, logs/metrics; Practical 4 & 10 mindset."),
]
for i, (t, d) in enumerate(outcomes):
    col = i % 2
    row = i // 2
    left = Inches(0.45 + col * 6.4)
    top = Inches(1.4 + row * 1.75)
    card(s, left, top, Inches(6.1), Inches(1.55), t, [d],
         accent=[SKY, TEAL, CORAL, GOLD, SKY, TEAL][i])

s = blank()
header(s, "Unit IV Syllabus — Roadmap", "Five syllabus topics  ·  each expanded with subtopics")
footer(s)
topics = [
    ("1", "Introduction to Kubernetes & its Architecture"),
    ("2", "Setting Up & Managing Kubernetes Clusters"),
    ("3", "Deploying Applications on Kubernetes"),
    ("4", "Kubernetes Networking & Load Balancing"),
    ("5", "Monitoring & Scaling Kubernetes Workloads"),
]
for i, (n, t) in enumerate(topics):
    top = Inches(1.35 + i * 1.0)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.85), WHITE)
    add_rect(s, Inches(0.5), top, Inches(0.9), Inches(0.85), [SKY, TEAL, CORAL, GOLD, GREEN][i])
    tb(s, Inches(0.65), top + Inches(0.25), Inches(0.6), Inches(0.4),
       [(n, 22, True, WHITE)], align=PP_ALIGN.CENTER)
    tb(s, Inches(1.65), top + Inches(0.25), Inches(10.8), Inches(0.4),
       [(t, 16, True, NAVY)])

# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 1 — Introduction & Architecture
# ═══════════════════════════════════════════════════════════════════════════

section_divider(1, "Introduction to Kubernetes & Architecture", [
    "Problem K8s solves · definition · containers vs VMs",
    "Control plane + worker nodes",
    "Must-know objects · desired vs actual state",
    "GKE Autopilot vs Standard",
])

s = blank()
header(s, "1.1  What problem does Kubernetes solve?", "One container is easy — hundreds across machines is not")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(1.4), "Without an orchestrator you must manually…",
     ["Decide which server runs which container · Restart crashed containers",
      "Scale before traffic spikes · Route users to healthy instances · Roll out without downtime"],
     accent=CORAL)
problems = [
    ("Placement", "Which node runs which container?"),
    ("Self-healing", "Restart when a container dies"),
    ("Scaling", "Add/remove replicas under load"),
    ("Networking", "Stable access as Pod IPs change"),
    ("Updates", "New version with near-zero downtime"),
    ("Ops at scale", "Hundreds of containers, many machines"),
]
for i, (t, d) in enumerate(problems):
    col, row = i % 3, i // 3
    card(s, Inches(0.45 + col * 4.25), Inches(3.0 + row * 1.9), Inches(4.05), Inches(1.7),
         t, [d], accent=[SKY, TEAL, GOLD, CORAL, GREEN, SKY][i])

s = blank()
header(s, "1.2  What is Kubernetes? (exam definition)", "Write this in long answers")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(2.2), "Definition",
     ["Kubernetes is an open-source container orchestration platform that automates",
      "the deployment, scaling, networking, and management of containerized",
      "applications across a cluster of machines.",
      "",
      "One line: “Kubernetes is the operating system for your containers at cluster scale.”"],
     accent=TEAL)
card(s, Inches(0.45), Inches(3.8), Inches(6.0), Inches(2.7), "Origin (viva)",
     ["Built by Google, inspired by Borg",
      "Now maintained by the CNCF community",
      "Often abbreviated K8s (K + 8 letters + s)"],
     accent=SKY)
card(s, Inches(6.7), Inches(3.8), Inches(6.15), Inches(2.7), "Skills transfer",
     ["GKE (GCP) · EKS (AWS) · AKS (Azure)",
      "Same core objects: Pod, Deployment, Service",
      "kubectl works across managed Kubernetes"],
     accent=GOLD)

s = blank()
header(s, "1.3  Container vs VM vs Kubernetes", "Quick orientation for exams")
footer(s)
card(s, Inches(0.45), Inches(1.4), Inches(4.05), Inches(4.8), "Virtual Machine",
     ["Full OS guest on a hypervisor",
      "Heavier · slower start",
      "Strong isolation",
      "",
      "Good for legacy / full OS needs"],
     accent=CORAL)
card(s, Inches(4.7), Inches(1.4), Inches(4.05), Inches(4.8), "Container",
     ["App + deps share host kernel",
      "Light · fast start",
      "Packaged as an image",
      "",
      "Docker / containerd run them"],
     accent=TEAL)
card(s, Inches(8.95), Inches(1.4), Inches(3.9), Inches(4.8), "Kubernetes",
     ["Schedules & manages many",
      "containers across nodes",
      "Self-healing · scale · LB",
      "",
      "Containers package; K8s operates"],
     accent=SKY)
banner(s, Inches(6.4), "Exam tip: Containers package the app; Kubernetes runs fleets of them.")

s = blank()
header(s, "1.4  Why Unit IV matters", "Course story arc")
footer(s)
arcs = [
    ("Unit I–II", "Design well · recover when infra fails", "Principles · DR · RTO/RPO"),
    ("Unit III", "Create infrastructure as code", "Terraform can create GKE too"),
    ("Unit IV", "Run microservices as containers at scale", "Orchestration · you are here"),
]
for i, (u, q, hint) in enumerate(arcs):
    left = Inches(0.45 + i * 4.25)
    add_rounded(s, left, Inches(1.5), Inches(4.05), Inches(4.5), WHITE)
    add_rect(s, left, Inches(1.5), Inches(4.05), Inches(0.7), [SKY, CORAL, TEAL][i])
    tb(s, left + Inches(0.2), Inches(1.65), Inches(3.6), Inches(0.4), [(u, 18, True, WHITE)])
    tb(s, left + Inches(0.25), Inches(2.5), Inches(3.55), Inches(1.8),
       [(q, 15, True, NAVY)])
    tb(s, left + Inches(0.25), Inches(4.6), Inches(3.55), Inches(0.8),
       [(hint, 13, False, MUTED)])

s = blank()
header(s, "1.5  Architecture — Control plane + worker nodes", "Draw this in 7-mark answers")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(1.15), "Cluster = control plane (brain) + worker nodes (muscle)",
     ["Users / kubectl talk to the API Server; controllers keep desired state true on nodes."],
     accent=TEAL)
cp = [
    ("API Server", "Front door for kubectl / controllers"),
    ("Scheduler", "Picks which node runs a Pod"),
    ("Controller Mgr", "Reconciles desired vs actual state"),
    ("etcd", "Distributed store of cluster state"),
]
for i, (t, d) in enumerate(cp):
    card(s, Inches(0.45 + i * 3.2), Inches(2.7), Inches(3.05), Inches(1.85),
         t, [d], accent=[SKY, TEAL, GOLD, CORAL][i])
wn = [
    ("kubelet", "Agent on each node — runs Pods as instructed"),
    ("kube-proxy", "Helps implement Service networking / LB rules"),
    ("Runtime", "containerd / Docker — actually starts containers"),
]
for i, (t, d) in enumerate(wn):
    card(s, Inches(0.45 + i * 4.25), Inches(4.75), Inches(4.05), Inches(1.85),
         t, [d], accent=[GREEN, SKY, TEAL][i])

s = blank()
header(s, "1.6  Must-know objects (vocabulary)", "Memorise names + one-line meaning")
footer(s)
objs = [
    ("Pod", "Smallest unit — one+ containers sharing net/volumes"),
    ("Deployment", "Desired replicas + rolling updates / rollbacks"),
    ("ReplicaSet", "Keeps N identical Pods (owned by Deployment)"),
    ("Service", "Stable VIP/DNS in front of changing Pods"),
    ("Namespace", "Virtual slice for isolation (dev / prod / team)"),
    ("ConfigMap / Secret", "Non-secret config / sensitive data for Pods"),
    ("Ingress", "HTTP(S) host/path routing to Services"),
    ("PV / PVC", "Durable storage claim for Pods needing disks"),
]
for i, (t, d) in enumerate(objs):
    col, row = i % 4, i // 4
    card(s, Inches(0.4 + col * 3.2), Inches(1.4 + row * 2.6), Inches(3.05), Inches(2.35),
         t, [d], accent=[SKY, TEAL, CORAL, GOLD][i % 4])

s = blank()
header(s, "1.7  Desired state vs actual state", "Same declarative idea as Terraform (Unit III)")
footer(s)
card(s, Inches(0.45), Inches(1.4), Inches(6.1), Inches(4.2), "You declare",
     ["“I want 3 replicas of hello-app image v2”",
      "",
      "Stored as Deployment / YAML in Git",
      "Controllers watch and reconcile",
      "",
      "If a Pod crashes → replacement appears"],
     accent=TEAL)
card(s, Inches(6.8), Inches(1.4), Inches(6.05), Inches(4.2), "Self-healing",
     ["Controllers continuously compare",
      "desired state vs actual state",
      "",
      "Create / replace / delete until match",
      "",
      "Exam phrase: “Desired state reconciled",
      "by controllers (self-healing)”"],
     accent=SKY)
banner(s, Inches(5.9), "Link Unit III: Terraform declares cloud resources; Kubernetes declares app replicas.")

s = blank()
header(s, "1.8  GKE — managed Kubernetes on Google Cloud", "Autopilot vs Standard (exam comparison)")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(1.25), "What is GKE?",
     ["Google Kubernetes Engine: Google runs the control plane; you deploy workloads.",
      "Open-source Kubernetes skills transfer to EKS / AKS / on-prem."],
     accent=TEAL)
card(s, Inches(0.45), Inches(2.85), Inches(6.1), Inches(3.6), "Autopilot",
     ["Google manages nodes",
      "Pay for Pod resource requests",
      "Opinionated safe defaults",
      "Best for most labs & many prod apps",
      "",
      "gcloud container clusters create-auto …"],
     accent=GREEN)
card(s, Inches(6.8), Inches(2.85), Inches(6.05), Inches(3.6), "Standard",
     ["You size / manage node pools",
      "Pay for node VMs",
      "More flexibility (GPUs, custom machines)",
      "Best when you need special hardware",
      "",
      "gcloud container clusters create …"],
     accent=GOLD)

# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 2 — Setting Up & Managing Clusters
# ═══════════════════════════════════════════════════════════════════════════

section_divider(2, "Setting Up & Managing Clusters", [
    "Big-picture setup steps · prerequisites",
    "Create Autopilot / Standard · get-credentials",
    "Release channels · namespaces · node ops",
    "Security baseline · IaC link (Unit III)",
])

s = blank()
header(s, "2.1  Cluster setup — big picture", "Practical 4 mindset")
footer(s)
steps = [
    ("1", "Project / billing / region"),
    ("2", "Enable APIs (Compute, GKE, Artifact Registry)"),
    ("3", "Create cluster (Autopilot or Standard)"),
    ("4", "get-credentials → kubeconfig"),
    ("5", "Verify: kubectl get nodes"),
    ("6", "Namespaces + pull permissions"),
]
for i, (n, t) in enumerate(steps):
    col, row = i % 3, i // 3
    left = Inches(0.45 + col * 4.25)
    top = Inches(1.4 + row * 2.6)
    add_rounded(s, left, top, Inches(4.05), Inches(2.35), WHITE)
    add_rect(s, left, top, Inches(4.05), Inches(0.65), [SKY, TEAL, CORAL, GOLD, GREEN, SKY][i])
    tb(s, left + Inches(0.2), top + Inches(0.15), Inches(3.6), Inches(0.4),
       [(f"Step {n}", 16, True, WHITE)])
    tb(s, left + Inches(0.25), top + Inches(0.95), Inches(3.55), Inches(1.1),
       [(t, 14, True, NAVY)])

s = blank()
header(s, "2.2  Prerequisites (lab / Practical 4)", "Tools & APIs before create")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(5.1), "You need",
     ["Google Cloud project with billing",
      "Tools: gcloud, kubectl, docker",
      "(Cloud Shell already has these)",
      "",
      "Enabled APIs:",
      "  · Compute Engine",
      "  · Kubernetes Engine",
      "  · Artifact Registry"],
     accent=SKY)
card(s, Inches(6.8), Inches(1.35), Inches(6.05), Inches(5.1), "Config sketch",
     ["export PROJECT_ID=your-project-id",
      "gcloud config set project $PROJECT_ID",
      "gcloud config set compute/region REGION",
      "",
      "Then create cluster + get-credentials",
      "",
      "Creation takes several minutes while",
      "control plane and nodes become healthy"],
     accent=TEAL)

s = blank()
header(s, "2.3  Creating a cluster on GKE", "Know the idea — exact flags may vary")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(4.5), "Autopilot (recommended for labs)",
     ["gcloud container clusters create-auto",
      "  hello-cluster --region=REGION",
      "",
      "Google manages nodes for you",
      "Focus on Deployments & Services"],
     accent=GREEN)
card(s, Inches(6.8), Inches(1.35), Inches(6.05), Inches(4.5), "Standard (more control)",
     ["gcloud container clusters create",
      "  hello-cluster --num-nodes=3",
      "  --region=REGION",
      "",
      "You choose node count / machine type",
      "Enable node auto-repair / auto-upgrade"],
     accent=GOLD)
banner(s, Inches(6.1), "Exam line: Managing a cluster = healthy control plane + nodes + correct credentials + namespaces.")

s = blank()
header(s, "2.4  Connecting kubectl", "get-credentials writes kubeconfig")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(2.0), "Connect",
     ["gcloud container clusters get-credentials hello-cluster --region=REGION",
      "",
      "Then: kubectl get nodes · kubectl cluster-info · kubectl get namespaces"],
     accent=TEAL)
card(s, Inches(0.45), Inches(3.6), Inches(6.1), Inches(2.9), "What happens?",
     ["Entry written into kubeconfig",
      "kubectl talks to that cluster’s API Server",
      "All later apply / get / logs use this context"],
     accent=SKY)
card(s, Inches(6.8), Inches(3.6), Inches(6.05), Inches(2.9), "Verify health",
     ["Nodes should show Ready",
      "If Pending Pods later → check capacity",
      "describe / events explain failures"],
     accent=CORAL)

s = blank()
header(s, "2.5  Release channels & upgrades", "Managing versions without surprise downtime")
footer(s)
channels = [
    ("Rapid", "Latest features — good for testing"),
    ("Regular", "Balance of new features & stability (common default)"),
    ("Stable", "Slowest, most conservative upgrades"),
]
for i, (t, d) in enumerate(channels):
    card(s, Inches(0.45 + i * 4.25), Inches(1.4), Inches(4.05), Inches(2.5),
         t, [d], accent=[CORAL, TEAL, GREEN][i])
card(s, Inches(0.45), Inches(4.15), Inches(12.4), Inches(2.4), "Management practices",
     ["Prefer regional clusters for higher control-plane availability",
      "Set maintenance windows so upgrades do not surprise exam / result day",
      "Test on non-prod first · keep kubectl client close to cluster version"],
     accent=GOLD)

s = blank()
header(s, "2.6  Namespaces & node operations", "Organise workloads · maintain nodes")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(5.1), "Namespaces",
     ["Virtual areas inside one cluster",
      "kubectl create namespace staging",
      "kubectl create namespace production",
      "",
      "kubectl get pods -n staging",
      "kubectl apply -f app.yaml -n production",
      "",
      "Separate teams, envs, or noisy tests"],
     accent=SKY)
card(s, Inches(6.8), Inches(1.35), Inches(6.05), Inches(5.1), "Node ops (Standard)",
     ["drain NODE — evacuate before maintenance",
      "cordon NODE — stop scheduling new Pods",
      "uncordon NODE — accept Pods again",
      "",
      "Autopilot: Google handles much of",
      "node repair / scaling for you"],
     accent=TEAL)

s = blank()
header(s, "2.7  Security baseline & IaC link", "Production habits + Unit III connection")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(5.1), "Security baseline",
     ["Least-privilege IAM for get-credentials",
      "Private nodes / authorised networks",
      "Workload Identity (avoid long-lived keys)",
      "NetworkPolicies for Pod-to-Pod limits",
      "Image scanning in Artifact Registry"],
     accent=CORAL)
card(s, Inches(6.8), Inches(1.35), Inches(6.05), Inches(5.1), "IaC link (Unit III)",
     ["Prod teams often create GKE with",
      "Terraform (google_container_cluster)",
      "",
      "Git is source of truth for cluster shape",
      "Then kubectl/Helm declare apps inside",
      "",
      "Two declarative layers: infra + workloads"],
     accent=TEAL)

# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 3 — Deploying Applications
# ═══════════════════════════════════════════════════════════════════════════

section_divider(3, "Deploying Applications on Kubernetes", [
    "Image → registry → Deployment → Service",
    "Pods, labels, declarative YAML",
    "Rolling updates & rollbacks",
    "ConfigMaps, Secrets, health probes",
])

s = blank()
header(s, "3.1  Deployment pipeline (container path)", "Standard path for Practical 4")
footer(s)
pipe = [
    ("Write app", "Code + Dockerfile"),
    ("Build image", "docker build …"),
    ("Push registry", "Artifact Registry"),
    ("Deployment", "replicas + template"),
    ("Expose Service", "ClusterIP / LB"),
    ("Verify", "pods · logs · IP"),
]
for i, (t, d) in enumerate(pipe):
    col, row = i % 3, i // 3
    left = Inches(0.45 + col * 4.25)
    top = Inches(1.4 + row * 2.6)
    add_rounded(s, left, top, Inches(4.05), Inches(2.35), WHITE)
    add_rect(s, left, top, Inches(4.05), Inches(0.65), [SKY, TEAL, CORAL, GOLD, GREEN, SKY][i])
    tb(s, left + Inches(0.2), top + Inches(0.15), Inches(3.6), Inches(0.4),
       [(t, 16, True, WHITE)])
    tb(s, left + Inches(0.25), top + Inches(0.95), Inches(3.55), Inches(1.1),
       [(d, 14, True, NAVY)])

s = blank()
header(s, "3.2  Container images & Artifact Registry", "GKE pulls images from a registry")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(1.3), "On GCP, Artifact Registry is the modern image store",
     ["Create repo → docker build/tag → configure-docker → push → grant pull permission"],
     accent=TEAL)
card(s, Inches(0.45), Inches(2.9), Inches(6.1), Inches(3.5), "Build & push (conceptual)",
     ["docker build -t REGION-docker.pkg.dev/",
      "  PROJECT/hello-repo/hello-app:v1 .",
      "",
      "gcloud auth configure-docker …",
      "docker push …/hello-app:v1"],
     accent=SKY)
card(s, Inches(6.8), Inches(2.9), Inches(6.05), Inches(3.5), "Permission tip",
     ["Grant node / GKE service account",
      "reader role on the repository",
      "",
      "Without pull rights → ImagePullBackOff",
      "Check: kubectl describe pod"],
     accent=CORAL)

s = blank()
header(s, "3.3  Pods & Deployments", "How apps are usually run")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(5.1), "Pod — smallest unit",
     ["Usually one main container",
      "(sometimes a helper sidecar)",
      "",
      "Pods are mortal — if a node dies,",
      "that Pod’s IP is gone",
      "",
      "Controllers recreate Pods;",
      "never treat a Pod IP as permanent"],
     accent=CORAL)
card(s, Inches(6.8), Inches(1.35), Inches(6.05), Inches(5.1), "Deployment — usual controller",
     ["Owns Pod template + replica count",
      "",
      "Imperative lab:",
      "kubectl create deployment hello-app",
      "  --image=…/hello-app:v1",
      "kubectl scale deployment hello-app",
      "  --replicas=3",
      "",
      "Prefer declarative YAML in Git"],
     accent=TEAL)

s = blank()
header(s, "3.4  Declarative Deployment YAML", "Show this pattern in exams / labs")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(5.15), "hello-deployment.yaml (shape)",
     ["apiVersion: apps/v1   kind: Deployment   metadata.name: hello-app",
      "spec.replicas: 3",
      "spec.selector.matchLabels.app: hello-app",
      "template.metadata.labels.app: hello-app",
      "template.spec.containers[0].image: …/hello-app:v1",
      "template.spec.containers[0].ports.containerPort: 8080",
      "",
      "kubectl apply -f hello-deployment.yaml"],
     accent=SKY)

s = blank()
header(s, "3.5  Labels, selectors & rolling updates", "Services find Pods by labels")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(3.4), "Labels & selectors",
     ["Labels = key/value tags (app=hello-app)",
      "Deployments & Services use selectors",
      "to find matching Pods",
      "",
      "Inconsistent labels → Service finds nothing"],
     accent=SKY)
card(s, Inches(6.8), Inches(1.35), Inches(6.05), Inches(3.4), "Rolling update",
     ["Build/push hello-app:v2",
      "kubectl set image deployment/hello-app …",
      "kubectl rollout status deployment/hello-app",
      "",
      "Old Pods gradually replaced; Service",
      "sends traffic only to Ready Pods"],
     accent=TEAL)
card(s, Inches(0.45), Inches(5.0), Inches(12.4), Inches(1.5), "Rollback if v2 is bad",
     ["kubectl rollout undo deployment/hello-app",
      "kubectl rollout history deployment/hello-app  ·  Fast reverse without SSH into machines"],
     accent=CORAL)

s = blank()
header(s, "3.6  ConfigMaps, Secrets & health probes", "Config without rebuild · Ready before traffic")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(4.05), Inches(5.1), "ConfigMap",
     ["Non-secret config",
      "Feature flags, env vars",
      "nginx snippets",
      "",
      "Inject without rebuilding image"],
     accent=SKY)
card(s, Inches(4.7), Inches(1.35), Inches(4.05), Inches(5.1), "Secret",
     ["Passwords, tokens, keys",
      "Base64 in etcd at rest",
      "Still protect with IAM",
      "",
      "Do not commit prod secrets to Git"],
     accent=CORAL)
card(s, Inches(8.95), Inches(1.35), Inches(3.9), Inches(5.1), "Probes",
     ["Liveness → restart if dead",
      "Readiness → remove from Service",
      "until ready",
      "Startup → slow booters",
      "",
      "Without readiness, users",
      "may hit booting Pods"],
     accent=TEAL)

s = blank()
header(s, "3.7  Common kubectl deploy commands", "Lab cheat sheet")
footer(s)
cmds = [
    ("kubectl apply -f FILE.yaml", "Create / update from manifest"),
    ("kubectl get pods,deploy,svc", "List key objects"),
    ("kubectl describe pod NAME", "Events — why Pending / CrashLoop"),
    ("kubectl logs POD [-c CONTAINER]", "Application logs"),
    ("kubectl exec -it POD -- /bin/sh", "Debug shell (use carefully)"),
    ("kubectl delete -f FILE.yaml", "Remove resources in file"),
]
for i, (c, d) in enumerate(cmds):
    top = Inches(1.3 + i * 0.9)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.75), WHITE)
    add_rect(s, Inches(0.5), top, Inches(0.12), Inches(0.75), [SKY, TEAL, CORAL, GOLD, GREEN, SKY][i])
    tb(s, Inches(0.85), top + Inches(0.1), Inches(7.5), Inches(0.5), [(c, 13, True, NAVY)])
    tb(s, Inches(8.5), top + Inches(0.18), Inches(4.0), Inches(0.45), [(d, 12, False, MUTED)])

# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 4 — Networking & Load Balancing
# ═══════════════════════════════════════════════════════════════════════════

section_divider(4, "Kubernetes Networking & Load Balancing", [
    "Why Pod IPs are not enough",
    "Service types (memorise)",
    "Ingress HTTP routing",
    "NetworkPolicy brief · CE vs GKE comparison",
])

s = blank()
header(s, "4.1  The core networking problem", "Pods change — clients need stability")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(1.6), "Problem",
     ["Pods get cluster IPs that change when Pods are recreated.",
      "Clients need a stable way to reach “the hello-app” as Pods come and go."],
     accent=CORAL)
layers = [
    ("Service", "Stable virtual IP + DNS → set of Pods"),
    ("Ingress", "HTTP(S) host/path routing → Services"),
    ("Cloud LB", "LoadBalancer type → cloud provisions external LB"),
]
for i, (t, d) in enumerate(layers):
    card(s, Inches(0.45 + i * 4.25), Inches(3.3), Inches(4.05), Inches(3.1),
         t, [d, "", "Must-say: Pod IPs are ephemeral —", "Services give stability"],
         accent=[SKY, TEAL, GOLD][i])

s = blank()
header(s, "4.2  Cluster DNS & Service types", "Memorise the four types")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(1.2), "DNS inside the cluster",
     ["hello-app.default.svc.cluster.local  ·  same namespace often just: hello-app"],
     accent=TEAL)
types = [
    ("ClusterIP", "Internal only (default)", "Backend microservices"),
    ("NodePort", "Node IP + high port", "Labs / simple exposure"),
    ("LoadBalancer", "External IP via cloud LB", "Public web / APIs"),
    ("ExternalName", "CNAME to external name", "Point outside the cluster"),
]
for i, (t, reach, use) in enumerate(types):
    card(s, Inches(0.4 + i * 3.2), Inches(2.85), Inches(3.05), Inches(3.6),
         t, [reach, "", f"Use: {use}"], accent=[SKY, TEAL, CORAL, GOLD][i])

s = blank()
header(s, "4.3  Expose publicly + port mapping", "Lab classic LoadBalancer path")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(2.5), "Expose (imperative lab)",
     ["kubectl expose deployment hello-app --name=hello-app-service",
      "  --type=LoadBalancer --port=80 --target-port=8080",
      "",
      "kubectl get service   → wait until EXTERNAL-IP leaves <pending>"],
     accent=SKY)
card(s, Inches(0.45), Inches(4.1), Inches(12.4), Inches(2.4), "Port mapping mental model",
     ["User → LoadBalancer port 80 → Service → Pod targetPort 8080 (container)",
      "",
      "Always match containerPort in the Pod with targetPort on the Service"],
     accent=TEAL)

s = blank()
header(s, "4.4  How load balancing works", "Ready endpoints join the pool")
footer(s)
card(s, Inches(0.45), Inches(1.4), Inches(6.1), Inches(4.8), "Service selects Pods by label",
     ["kube-proxy / dataplane distributes",
      "connections across Ready endpoints",
      "",
      "When HPA adds Pods, they join the pool",
      "once Ready (readiness probes matter)",
      "",
      "Unready Pods do not receive traffic"],
     accent=TEAL)
card(s, Inches(6.8), Inches(1.4), Inches(6.05), Inches(4.8), "Real-life parallel",
     ["UPI-style stack:",
      "Ingress terminates TLS at the edge",
      "Routes /pay to payments ClusterIP",
      "payments Pods stay private",
      "",
      "Only Ingress / LB edge is public"],
     accent=GOLD)

s = blank()
header(s, "4.5  Ingress — HTTP(S) routing", "One entry · many services")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(1.5), "Example routes",
     ["api.example.com/orders → orders-service",
      "api.example.com/pay → payments-service"],
     accent=TEAL)
card(s, Inches(0.45), Inches(3.1), Inches(6.1), Inches(3.4), "Benefits",
     ["One external entry, many services",
      "Host / path based routing",
      "TLS termination with certificates",
      "(production pattern)"],
     accent=SKY)
card(s, Inches(6.8), Inches(3.1), Inches(6.05), Inches(3.4), "Exam contrast",
     ["Service LoadBalancer → often one",
      "LB per service (simple)",
      "",
      "Ingress → shared HTTP routing layer",
      "(microservices friendly)"],
     accent=CORAL)

s = blank()
header(s, "4.6  NetworkPolicy & CE vs GKE networking", "Brief advanced note + Unit II lab memory")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(5.1), "NetworkPolicy (brief)",
     ["By default many clusters allow",
      "wide Pod-to-Pod traffic",
      "",
      "NetworkPolicy can restrict who may talk",
      "(zero-trust style)",
      "",
      "Important for multi-tenant & payments"],
     accent=CORAL)
card(s, Inches(6.8), Inches(1.35), Inches(6.05), Inches(5.1), "Compute Engine vs GKE",
     ["CE: firewall rules + VM IPs +",
      "HTTP(S) LB + MIG",
      "",
      "GKE: Services / Ingress + Pod IPs +",
      "cloud LB via Service type",
      "",
      "Same goal (reach healthy backends);",
      "different abstraction level"],
     accent=TEAL)

# ═══════════════════════════════════════════════════════════════════════════
# TOPIC 5 — Monitoring & Scaling
# ═══════════════════════════════════════════════════════════════════════════

section_divider(5, "Monitoring & Scaling Workloads", [
    "Why monitor · metrics / logs / traces",
    "HPA (Horizontal Pod Autoscaler)",
    "Cluster / node scaling · VPA brief",
    "Alerting · logging · peak events",
])

s = blank()
header(s, "5.1  Why monitoring matters", "Orchestration hides machines — you must watch")
footer(s)
watch = [
    ("Restarts", "CrashLoopBackOff / restart storms"),
    ("Saturation", "CPU & memory pressure"),
    ("Latency", "Request latency & error rate"),
    ("Nodes", "NotReady · disk pressure"),
    ("Rollouts", "Deployment rollout health"),
    ("Capacity", "Pods Pending too long"),
]
for i, (t, d) in enumerate(watch):
    col, row = i % 3, i // 3
    card(s, Inches(0.45 + col * 4.25), Inches(1.4 + row * 2.6), Inches(4.05), Inches(2.35),
         t, [d, "", "Practical 10 focus:", "monitor & log K8s workloads"],
         accent=[SKY, TEAL, CORAL, GOLD, GREEN, SKY][i])

s = blank()
header(s, "5.2  Signals — metrics, logs, traces", "Cloud Operations on GKE")
footer(s)
card(s, Inches(0.45), Inches(1.4), Inches(4.05), Inches(5.0), "Metrics",
     ["CPU, memory, RPS, latency",
      "Feed HPA and dashboards",
      "",
      "kubectl top nodes / pods",
      "(needs metrics-server)"],
     accent=SKY)
card(s, Inches(4.7), Inches(1.4), Inches(4.05), Inches(5.0), "Logs",
     ["stdout / stderr from containers",
      "kubectl logs · Cloud Logging",
      "",
      "12-factor: log to stdout",
      "Filter by labels / namespaces"],
     accent=TEAL)
card(s, Inches(8.95), Inches(1.4), Inches(3.9), Inches(5.0), "Traces",
     ["Request path across",
      "microservices (APM)",
      "",
      "Advanced but useful for",
      "payment / checkout flows"],
     accent=GOLD)

s = blank()
header(s, "5.3  Useful observability commands", "Diagnose before you guess")
footer(s)
obs = [
    ("kubectl get pods -o wide", "Where Pods run · IPs · nodes"),
    ("kubectl describe pod NAME", "Events explain Pending / OOMKilled"),
    ("kubectl top nodes / pods", "CPU & memory snapshot"),
    ("kubectl logs deploy/hello-app", "App output / errors"),
    ("kubectl get events --sort-by=…", "Recent cluster events"),
]
for i, (c, d) in enumerate(obs):
    top = Inches(1.35 + i * 1.05)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.9), WHITE)
    add_rect(s, Inches(0.5), top, Inches(0.12), Inches(0.9), [SKY, TEAL, CORAL, GOLD, GREEN][i])
    tb(s, Inches(0.85), top + Inches(0.12), Inches(7.8), Inches(0.3), [(c, 13, True, NAVY)])
    tb(s, Inches(0.85), top + Inches(0.45), Inches(11.5), Inches(0.35), [(d, 12, False, MUTED)])

s = blank()
header(s, "5.4  Horizontal Pod Autoscaler (HPA)", "Scale Pods — not VMs (by default)")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(1.6), "Command",
     ["kubectl autoscale deployment hello-app --cpu-percent=80 --min=1 --max=5",
      "If average CPU > target → add Pods (up to max) · load drops → remove (down to min)"],
     accent=TEAL)
card(s, Inches(0.45), Inches(3.2), Inches(6.1), Inches(3.3), "Behaviour",
     ["Needs metrics (metrics-server /",
      "Cloud Monitoring)",
      "",
      "Set resource requests/limits so",
      "scheduler & autoscaler know needs",
      "",
      "Analogy: MIG scales VMs; HPA scales Pods"],
     accent=SKY)
card(s, Inches(6.8), Inches(3.2), Inches(6.05), Inches(3.3), "Peak events",
     ["Result day / festival sale:",
      "raise HPA max; ensure cluster",
      "can add nodes; watch error rate",
      "",
      "After peak: let HPA scale down",
      "(FinOps — avoid idle spend)"],
     accent=GOLD)

s = blank()
header(s, "5.5  Cluster / node scaling & VPA", "When Pods are Pending — capacity problem")
footer(s)
card(s, Inches(0.45), Inches(1.35), Inches(6.1), Inches(5.1), "Node / cluster scaling",
     ["Autopilot: Google adds capacity",
      "based on Pod requests",
      "",
      "Standard: Cluster Autoscaler /",
      "node auto-provisioning can add nodes",
      "",
      "HPA scales Pods;",
      "cluster autoscaler scales nodes"],
     accent=TEAL)
card(s, Inches(6.8), Inches(1.35), Inches(6.05), Inches(5.1), "Vertical Pod Autoscaler (brief)",
     ["VPA suggests or applies higher/lower",
      "CPU/memory requests",
      "",
      "Often combined carefully with HPA",
      "(advanced — know the idea)",
      "",
      "Exam: name HPA first; mention VPA",
      "as vertical alternative"],
     accent=SKY)

s = blank()
header(s, "5.6  Alerting & logging best practices", "Page on symptoms users feel")
footer(s)
alerts = [
    "CrashLoopBackOff / high restart rate",
    "Elevated 5xx at Ingress / Service",
    "Pod Pending too long (capacity)",
    "Node NotReady",
    "Apps log to stdout/stderr → Cloud Logging",
    "Filter by labels / namespaces per team",
]
for i, a in enumerate(alerts):
    top = Inches(1.35 + i * 0.85)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.7), WHITE)
    add_rect(s, Inches(0.5), top, Inches(0.12), Inches(0.7), [SKY, TEAL, CORAL, GOLD, GREEN, SKY][i])
    tb(s, Inches(0.85), top + Inches(0.18), Inches(11.7), Inches(0.4), [(a, 13, False, INK)])

# ═══════════════════════════════════════════════════════════════════════════
# CASE STUDIES
# ═══════════════════════════════════════════════════════════════════════════

s = blank()
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
add_rect(s, Inches(0), Inches(0), Inches(0.25), SLIDE_H, GOLD)
tb(s, Inches(0.8), Inches(2.3), Inches(11), Inches(0.4),
   [("CASE STUDIES", 14, True, GOLD)])
tb(s, Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.0),
   [("Real-life examples for 7-mark answers", 32, True, WHITE)])
tb(s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.8),
   [("Quote one case + object (Deployment/Service/HPA) + teaching point", 14, False, SOFT)])
footer(s)

cases = [
    ("Campus result portal on GKE",
     "Container + Deployment + HPA min=2 max=20 + LoadBalancer/Ingress + 5xx alerts.",
     "Scale Pods automatically; rolling update content without overnight downtime."),
    ("Food delivery microservices",
     "Many Deployments; ClusterIP between services; Ingress at edge; HPA per service; namespaces.",
     "Networking + independent replica counts fit microservices."),
    ("UPI / payments API",
     "Regional GKE; readiness probes; staged rolling update; fast rollout undo; NetworkPolicy.",
     "Orchestration + probes + rollback protect trust-critical traffic."),
    ("E-commerce festival sale",
     "Same manifests; staging namespace with smaller HPA max; promote image tag after soak.",
     "Declarative deploys make environment parity practical."),
    ("Batch analytics on Spot",
     "Jobs/CronJobs on Spot nodes; checkpoint to object storage; notify on failure, not every preempt.",
     "Not every workload needs Hot Standby replicas — choose controllers wisely."),
    ("Bad image → CrashLoop",
     "describe/logs diagnose; rollout undo; fix ConfigMap; add readiness probe.",
     "Monitoring commands + rollback are daily Kubernetes operations."),
]
for batch in (cases[:3], cases[3:]):
    s = blank()
    header(s, "Real-life cases", "Use these stories in long answers & viva")
    footer(s)
    for i, (title, approach, point) in enumerate(batch):
        top = Inches(1.35 + i * 1.8)
        add_rounded(s, Inches(0.45), top, Inches(12.4), Inches(1.6), WHITE)
        add_rect(s, Inches(0.45), top, Inches(0.12), Inches(1.6), [SKY, TEAL, CORAL][i])
        tb(s, Inches(0.8), top + Inches(0.15), Inches(11.8), Inches(0.35),
           [(title, 15, True, NAVY)])
        tb(s, Inches(0.8), top + Inches(0.55), Inches(11.8), Inches(0.4),
           [(approach, 12, False, INK)])
        tb(s, Inches(0.8), top + Inches(1.05), Inches(11.8), Inches(0.35),
           [(f"Teaching point: {point}", 12, True, MUTED)])

# ═══════════════════════════════════════════════════════════════════════════
# EXAM + CLOSE
# ═══════════════════════════════════════════════════════════════════════════

s = blank()
header(s, "Exam Rapid Revision — Sheet A", "Definitions & architecture")
footer(s)
lines = [
    "Kubernetes = container orchestration across a cluster (deploy, scale, network, heal)",
    "Control plane: API Server · Scheduler · Controller Manager · etcd",
    "Node: kubelet · kube-proxy · container runtime",
    "Pod = smallest unit · Deployment = replicas + rolling updates · Service = stable endpoint",
    "Desired state reconciled by controllers (self-healing)",
    "GKE Autopilot: Google manages nodes, pay per Pod requests · Standard: you manage nodes",
    "Pod IPs ephemeral — Services / Ingress give stability",
]
for i, line in enumerate(lines):
    top = Inches(1.3 + i * 0.75)
    add_rect(s, Inches(0.5), top + Inches(0.1), Inches(0.16), Inches(0.16),
             SKY if i % 2 == 0 else TEAL)
    tb(s, Inches(0.9), top, Inches(11.8), Inches(0.65), [(line, 13, False, INK)])

s = blank()
header(s, "Exam Rapid Revision — Sheet B", "Deploy · network · scale · must-say phrases")
footer(s)
lines = [
    "Deploy path: Build image → Registry → Deployment → Service → Verify → (HPA)",
    "Rolling update + kubectl rollout undo = safer releases",
    "Service types: ClusterIP · NodePort · LoadBalancer · ExternalName",
    "Ingress = HTTP(S) host/path routing (and TLS) to many Services",
    "HPA scales Pods; cluster autoscaler scales nodes",
    "Must-say: Desired state · ephemeral Pod IPs · Ready before traffic · Practical 4 & 10",
    "Link Unit III: Terraform creates GKE; Kubernetes YAML declares apps inside",
]
for i, line in enumerate(lines):
    top = Inches(1.3 + i * 0.75)
    add_rect(s, Inches(0.5), top + Inches(0.1), Inches(0.16), Inches(0.16),
             CORAL if i % 2 == 0 else GOLD)
    tb(s, Inches(0.9), top, Inches(11.8), Inches(0.65), [(line, 13, False, INK)])

s = blank()
header(s, "Important exam questions map", "Q → which topic to revise")
footer(s)
qs = [
    ("Q1", "What is Kubernetes? Why needed?", "Topic 1"),
    ("Q2", "Explain architecture (control plane & nodes)", "Topic 1"),
    ("Q3", "Pod, Deployment, Service definitions", "Topics 1 + 3"),
    ("Q4", "Set up & connect to a GKE cluster", "Topic 2"),
    ("Q5", "Steps to deploy a containerized app", "Topic 3"),
    ("Q6", "Rolling update & rollback", "Topic 3"),
    ("Q7", "Service types · when use Ingress?", "Topic 4"),
    ("Q8", "HPA + monitoring · Autopilot vs Standard", "Topics 1 + 5"),
]
for i, (q, prompt, where) in enumerate(qs):
    col, row = i % 2, i // 2
    left = Inches(0.45 + col * 6.4)
    top = Inches(1.35 + row * 1.35)
    add_rounded(s, left, top, Inches(6.15), Inches(1.2), WHITE)
    tb(s, left + Inches(0.2), top + Inches(0.15), Inches(1.0), Inches(0.35),
       [(q, 14, True, NAVY)])
    tb(s, left + Inches(1.2), top + Inches(0.18), Inches(3.4), Inches(0.8),
       [(prompt, 12, False, INK)])
    tb(s, left + Inches(4.5), top + Inches(0.4), Inches(1.5), Inches(0.4),
       [(where, 11, True, TEAL)])

s = blank()
header(s, "Key Takeaways", "Operate containers like a fleet — not one-by-one SSH")
footer(s)
takes = [
    ("Declare desired state", "Controllers reconcile — self-healing is the product."),
    ("Never trust a Pod IP", "Expose via Services / Ingress; labels must stay consistent."),
    ("Ship with rollouts", "Rolling update + undo beat overnight manual upgrades."),
    ("Scale with signals", "HPA needs metrics; watch Pending Pods for capacity."),
    ("Two layers with IaC", "Terraform shapes the cluster; YAML shapes the apps (Unit III)."),
]
for i, (t, d) in enumerate(takes):
    top = Inches(1.3 + i * 1.05)
    add_rounded(s, Inches(0.5), top, Inches(12.3), Inches(0.9), WHITE)
    add_rect(s, Inches(0.5), top, Inches(0.12), Inches(0.9), [SKY, TEAL, CORAL, GOLD, GREEN][i])
    tb(s, Inches(0.9), top + Inches(0.12), Inches(11.5), Inches(0.3), [(t, 14, True, NAVY)])
    tb(s, Inches(0.9), top + Inches(0.45), Inches(11.5), Inches(0.35), [(d, 12, False, MUTED)])

s = blank()
add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.15), TEAL)
tb(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.8), [("Thank you", 44, True, WHITE)])
tb(s, Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.5),
   [("Questions · Pick a case · Map it to Pod / Deployment / Service / HPA", 16, False, SKY_LIGHT)])
tb(s, Inches(0.8), Inches(4.2), Inches(11.5), Inches(1.6),
   [("Unit IV — Kubernetes & Container Orchestration", 14, False, SOFT),
    ("Pair with: notes/unit-4/CC-II-Unit-IV-Kubernetes-Container-Orchestration.txt", 14, False, SOFT),
    ("Lab links: Practical 4 (setup + deploy + scale) · Practical 10 (monitor + log)", 14, False, SOFT),
    ("FIT & CS, Parul University", 14, False, MUTED)])
PAGE += 1

out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "CC-II-Unit-IV-Kubernetes-Container-Orchestration.pptx")
prs.save(out)
print(f"Saved {out}")
print(f"Slides: {len(prs.slides)}")
